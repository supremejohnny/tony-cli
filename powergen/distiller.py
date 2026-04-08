from __future__ import annotations

import hashlib
import json
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from .mock_client import LLMClient
    from .workspace import WorkspaceContext

from .prompts_distill import (
    distill_generic_system_prompt,
    distill_generic_user_prompt,
    distill_pptx_system_prompt,
    distill_pptx_user_prompt,
    vision_describe_system_prompt,
    vision_describe_user_prompt,
)

# Extensions handled by each extraction path
_PPTX_EXTS = {".pptx"}
_MARKITDOWN_EXTS = {".pdf", ".docx"}
_PLAINTEXT_EXTS = {".md", ".txt"}
_ALL_SUPPORTED = _PPTX_EXTS | _MARKITDOWN_EXTS | _PLAINTEXT_EXTS

# Guard: truncate input if it would blow past a reasonable context budget
_MAX_INPUT_CHARS = 40_000

# Directory names that are build/packaging artefacts — skip files inside them
_SKIP_DIRS = {
    "__pycache__", "node_modules", "dist", "build", "site-packages",
    ".egg-info",   # matched as suffix below
}


@dataclass
class SlideText:
    index: int   # 1-based
    title: str
    body: str


@dataclass
class SlideImage:
    slide_index: int   # 1-based, matches SlideText.index
    blob: bytes
    media_type: str    # "image/jpeg" | "image/png" | "image/gif" | "image/webp"


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_pptx_slides(path: Path) -> list[SlideText]:
    """Extract per-slide title + body text using python-pptx."""
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(str(path))
    slides: list[SlideText] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Identify title placeholder: placeholder index 0 is the title
            try:
                if shape.placeholder_format is not None and shape.placeholder_format.idx == 0:
                    title = text
                    continue
            except Exception:
                pass
            body_parts.append(text)
        slides.append(SlideText(index=i, title=title, body="\n".join(body_parts)))
    return slides


def _detect_media_type(blob: bytes) -> str:
    """Detect image media type from magic bytes."""
    if blob[:4] == b"\x89PNG":
        return "image/png"
    if blob[:2] == b"\xff\xd8":
        return "image/jpeg"
    if blob[:4] in (b"GIF8", b"GIF9"):
        return "image/gif"
    if blob[:4] == b"RIFF" and blob[8:12] == b"WEBP":
        return "image/webp"
    return "image/png"  # safe fallback


def _is_decorative(shape: object, slide_w: int, slide_h: int) -> bool:
    """Return True if the shape should be skipped (too small or in the margin)."""
    area = getattr(shape, "width", 0) * getattr(shape, "height", 0)
    if area < 0.05 * slide_w * slide_h:
        return True
    cx = getattr(shape, "left", 0) + getattr(shape, "width", 0) / 2
    cy = getattr(shape, "top", 0) + getattr(shape, "height", 0) / 2
    if (cx < slide_w * 0.10 or cx > slide_w * 0.90
            or cy < slide_h * 0.10 or cy > slide_h * 0.90):
        return True
    return False


def _iter_shapes(shapes: object):
    """Recursively yield all shapes, descending into group shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import]
    for shape in shapes:  # type: ignore[union-attr]
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def extract_pptx_slides_with_images(
    path: Path,
) -> tuple[list[SlideText], list[SlideImage]]:
    """Extract per-slide text and embedded images from a PPTX file.

    Images are detected by attempting shape.image.blob (works for both
    PICTURE and PLACEHOLDER shapes that contain images). Group shapes are
    traversed recursively.

    Filters: area >= 5% of slide AND centre not in outer 10% margin.
    Blobs > 3.5 MB are skipped (Anthropic API per-image size limit).
    """
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(str(path))
    slide_w: int = prs.slide_width
    slide_h: int = prs.slide_height

    slides: list[SlideText] = []
    images: list[SlideImage] = []

    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body_parts: list[str] = []
        for shape in _iter_shapes(slide.shapes):
            # --- text extraction ---
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    try:
                        if (shape.placeholder_format is not None
                                and shape.placeholder_format.idx == 0):
                            title = text
                            continue
                    except Exception:
                        pass
                    body_parts.append(text)

            # --- image extraction ---
            # Use image.blob access rather than shape_type check:
            # images in PLACEHOLDER shapes have type=14, not PICTURE=13.
            try:
                blob = shape.image.blob
            except (AttributeError, ValueError):
                blob = None  # shape has no image — expected, skip
            except Exception as exc:
                print(f"  [warn] slide {i} shape '{shape.name}': unexpected image read error: {exc}")
                blob = None

            if blob is not None:
                if not _is_decorative(shape, slide_w, slide_h) and len(blob) <= 3_500_000:
                    images.append(SlideImage(
                        slide_index=i,
                        blob=blob,
                        media_type=_detect_media_type(blob),
                    ))

        slides.append(SlideText(index=i, title=title, body="\n".join(body_parts)))

    return slides, images


def format_slides_for_prompt(slides: list[SlideText]) -> str:
    """Format slide list into a readable block for the LLM prompt."""
    parts: list[str] = []
    for s in slides:
        header = f"=== Slide {s.index}: {s.title} ===" if s.title else f"=== Slide {s.index} ==="
        content = s.body if s.body else "(no text)"
        parts.append(f"{header}\n{content}")
    return "\n\n".join(parts)


def extract_generic_text(path: Path) -> str:
    """Extract text from PDF/DOCX (markitdown) or MD/TXT (plain read)."""
    ext = path.suffix.lower()
    if ext in _MARKITDOWN_EXTS:
        try:
            from markitdown import MarkItDown  # type: ignore[import]
            md = MarkItDown()
            result = md.convert(str(path))
            return result.text_content
        except ImportError as e:
            raise RuntimeError(
                f"markitdown is required to process {ext} files. "
                "Install it with: pip install 'markitdown[pptx]'"
            ) from e
    else:
        return path.read_text(encoding="utf-8", errors="replace")


# ---------------------------------------------------------------------------
# Hashing
# ---------------------------------------------------------------------------

def compute_file_hash(path: Path) -> str:
    """Return 'sha256:<hex>' for the file at *path*."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return f"sha256:{h.hexdigest()}"


# ---------------------------------------------------------------------------
# Response parsing
# ---------------------------------------------------------------------------

def _parse_distill_response(raw: str, label: str) -> dict:
    """Extract JSON from LLM response using a 3-step fallback."""
    text = raw.strip()

    # Step 1: direct parse
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # Step 2: extract from ```json ... ``` or ``` ... ``` fence
    match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", text, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(1))
        except json.JSONDecodeError:
            pass

    # Step 3: first bare { ... } block
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(0))
        except json.JSONDecodeError:
            pass

    raise ValueError(
        f"Could not parse {label} distill response as JSON.\n"
        f"Response preview:\n{raw[:500]}"
    )


# ---------------------------------------------------------------------------
# Vision — describe images embedded in PPTX slides
# ---------------------------------------------------------------------------

def _describe_slide_images(
    images: list[SlideImage],
    client: "LLMClient",
) -> dict[int, list[str]]:
    """Call generate_vision() once per slide that has images.

    Returns a mapping of slide_index → list of '[Visual: ...]' strings.
    All images from the same slide are batched into a single API call.
    """
    import base64

    # Group images by slide
    by_slide: dict[int, list[SlideImage]] = {}
    for img in images:
        by_slide.setdefault(img.slide_index, []).append(img)

    system = vision_describe_system_prompt()
    result: dict[int, list[str]] = {}

    for slide_idx, slide_images in sorted(by_slide.items()):
        image_blocks: list[dict] = [
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": img.media_type,
                    "data": base64.standard_b64encode(img.blob).decode("ascii"),
                },
            }
            for img in slide_images
        ]
        user = vision_describe_user_prompt(slide_idx)
        try:
            raw = client.generate_vision(system, image_blocks, user).strip()
        except Exception as exc:
            raw = f"[Vision error: {exc}]"

        if not raw or raw.lower() == "decorative":
            result[slide_idx] = []
            continue

        if not raw.startswith("[Visual:"):
            raw = f"[Visual: {raw}]"
        result[slide_idx] = [raw]

    return result


# ---------------------------------------------------------------------------
# Local combined_text injection for PPTX
# ---------------------------------------------------------------------------

def _build_slide_index(slides: list[SlideText]) -> dict[int, SlideText]:
    """Return a 1-based index mapping slide number → SlideText."""
    return {s.index: s for s in slides}


def _inject_combined_text(
    data: dict,
    slide_index: dict[int, SlideText],
    vision_map: dict[int, list[str]] | None = None,
) -> None:
    """Inject combined_text (and has_images) into each PPTX chunk.

    The LLM does not output combined_text for PPTX — we assemble it locally
    from the already-extracted SlideText list, guaranteeing verbatim content.
    Vision descriptions from vision_map are appended after the slide text.
    Both fields are inserted right after 'titles' so the JSON reads naturally.
    """
    for chunk in data.get("chunks", []):
        slide_range = chunk.get("slide_range")
        if not slide_range:
            continue
        start = int(slide_range[0])
        end = int(slide_range[-1])  # works for both [N] and [N, M]
        parts: list[str] = []
        for idx in range(start, end + 1):
            slide = slide_index.get(idx)
            if slide is None:
                continue
            header = f"=== Slide {slide.index}: {slide.title} ===" if slide.title else f"=== Slide {slide.index} ==="
            body = slide.body if slide.body else "(no text)"
            parts.append(f"{header}\n{body}")
        combined = "\n\n".join(parts)

        # Append vision descriptions and determine has_images
        has_images = False
        if vision_map:
            for idx in range(start, end + 1):
                descs = vision_map.get(idx, [])
                if descs:
                    has_images = True
                    combined += "\n" + "\n".join(descs)

        # Insert combined_text and has_images after 'titles' by rebuilding the chunk dict
        new_chunk: dict = {}
        for key, val in chunk.items():
            new_chunk[key] = val
            if key == "titles":
                new_chunk["combined_text"] = combined
                new_chunk["has_images"] = has_images
        chunk.clear()
        chunk.update(new_chunk)


# ---------------------------------------------------------------------------
# Per-file distillation
# ---------------------------------------------------------------------------

def _distill_one(path: Path, client: "LLMClient", enable_vision: bool = True) -> dict:
    """Run a single LLM distill call for *path*. Returns the parsed JSON dict."""
    ext = path.suffix.lower()

    if ext in _PPTX_EXTS:
        if enable_vision:
            slides, images = extract_pptx_slides_with_images(path)
        else:
            slides = extract_pptx_slides(path)
            images = []
        slide_text = format_slides_for_prompt(slides)
        if len(slide_text) > _MAX_INPUT_CHARS:
            slide_text = slide_text[:_MAX_INPUT_CHARS] + "\n\n[... truncated — file too large ...]"
        system = distill_pptx_system_prompt()
        user = distill_pptx_user_prompt(path.name, slide_text)
        raw = client.generate(system, user)
        data = _parse_distill_response(raw, path.name)
        # Build vision map from images (empty if vision disabled or no images found)
        vision_map = _describe_slide_images(images, client) if images else None
        # Inject combined_text locally — no LLM output tokens spent on verbatim text
        _inject_combined_text(data, _build_slide_index(slides), vision_map)
    else:
        content = extract_generic_text(path)
        if len(content) > _MAX_INPUT_CHARS:
            content = content[:_MAX_INPUT_CHARS] + "\n\n[... truncated — file too large ...]"
        system = distill_generic_system_prompt()
        user = distill_generic_user_prompt(path.name, content)
        raw = client.generate(system, user)
        data = _parse_distill_response(raw, path.name)

    # Stamp in caller-side fields that LLM left as ""
    now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    file_hash = compute_file_hash(path)
    if "source" not in data:
        data["source"] = {}
    data["source"]["file_name"] = path.name
    data["source"]["file_hash"] = file_hash
    data["source"]["distilled_at"] = now

    return data


# ---------------------------------------------------------------------------
# Hash-based skip check
# ---------------------------------------------------------------------------

def _is_current(distill_path: Path, file_hash: str) -> bool:
    """Return True if *distill_path* exists and records the same hash."""
    if not distill_path.exists():
        return False
    try:
        existing = json.loads(distill_path.read_text(encoding="utf-8"))
        stored_hash = existing.get("source", {}).get("file_hash", "")
        return stored_hash == file_hash
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Index update
# ---------------------------------------------------------------------------

def _update_index(distill_dir: Path, entry: dict) -> None:
    """Insert or replace the entry for this source file in _index.json."""
    index_path = distill_dir / "_index.json"
    if index_path.exists():
        try:
            index: list[dict] = json.loads(index_path.read_text(encoding="utf-8"))
        except Exception:
            index = []
    else:
        index = []

    source_file = entry.get("source_file", "")
    index = [e for e in index if e.get("source_file") != source_file]
    index.append(entry)
    index_path.write_text(json.dumps(index, indent=2, ensure_ascii=False), encoding="utf-8")


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def _distill_stem(path: Path, cwd: Path) -> str:
    """Return a collision-safe stem for the distill filename.

    Root-level files → plain stem ("README").
    Nested files → path parts joined with "_" ("rust_README").
    """
    try:
        rel = path.relative_to(cwd)
    except ValueError:
        return path.stem
    parts = list(rel.parts)  # e.g. ["rust", "README.md"]
    parts[-1] = Path(parts[-1]).stem   # strip extension from last part
    if len(parts) == 1:
        return parts[0]
    return "_".join(parts)


def run_distill(
    workspace: "WorkspaceContext",
    client: "LLMClient",
    distill_dir: Path,
    force: bool = False,
    enable_vision: bool = True,
) -> None:
    """Distill all supported files in *workspace*, writing results to *distill_dir*."""
    cwd = Path.cwd()

    # Collect paths to process
    candidates: list[Path] = []
    for ti in workspace.templates:
        if ti.path.suffix.lower() in _PPTX_EXTS:
            candidates.append(ti.path)
    for wf in workspace.content_files:
        if wf.path.suffix.lower() in (_MARKITDOWN_EXTS | _PLAINTEXT_EXTS):
            candidates.append(wf.path)

    # Filter out build/packaging artefact directories
    def _in_skip_dir(p: Path) -> bool:
        for part in p.parts:
            if part in _SKIP_DIRS or part.endswith(".egg-info"):
                return True
        return False

    candidates = [p for p in candidates if not _in_skip_dir(p)]

    if not candidates:
        print("No supported files found to distill (.pptx, .pdf, .docx, .md, .txt).")
        return

    print(f"Scanning {len(candidates)} file(s)...")

    for i, path in enumerate(candidates, 1):
        distill_filename = _distill_stem(path, cwd) + ".distill.json"
        distill_path = distill_dir / distill_filename
        prefix = f"[{i}/{len(candidates)}] {path.name}"

        if not force:
            file_hash = compute_file_hash(path)
            if _is_current(distill_path, file_hash):
                print(f"{prefix} -- unchanged, skipping.")
                continue

        print(f"{prefix} -> distilling...")
        try:
            data = _distill_one(path, client, enable_vision=enable_vision)
        except Exception as exc:
            print(f"  Error: {exc}")
            continue

        distill_path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        print(f"  Saved: {distill_path.name}")

        # Update collection index — use relative path as key to avoid same-name collisions
        try:
            rel_source = str(path.relative_to(cwd))
        except ValueError:
            rel_source = path.name
        index_entry = {
            "source_file": rel_source,
            "distill_file": distill_filename,
            "file_hash": data.get("source", {}).get("file_hash", ""),
            "distilled_at": data.get("source", {}).get("distilled_at", ""),
            "global_summary": data.get("global_summary", ""),
            "main_topics": data.get("main_topics", []),
        }
        _update_index(distill_dir, index_entry)

    index_path = distill_dir / "_index.json"
    print(f"Index: {index_path}")
    print("Done.")
