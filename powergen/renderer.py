from __future__ import annotations

import re
from pathlib import Path

from .models import PresentationSpec, SlideSpec

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Characters not allowed in Windows filenames (conservative cross-platform set)
_UNSAFE_CHARS = re.compile(r'[<>:"/\\|?*\x00-\x1f]')


def _safe_filename(title: str) -> str:
    """Convert a presentation title to a safe .pptx filename."""
    name = _UNSAFE_CHARS.sub("-", title).strip(". ")
    name = re.sub(r"-{2,}", "-", name)  # collapse multiple dashes
    return (name or "presentation") + ".pptx"


def render_presentation(
    spec: PresentationSpec,
    output_path: Path | None = None,
    template_path: Path | None = None,
) -> Path:
    """Render a PresentationSpec to a .pptx file.

    Template mode: open the template, remove all existing slides (correctly,
    via drop_rel so the ZIP parts are actually released), then add spec slides.
    No-template mode: create a blank Presentation and add spec slides.
    Returns the resolved output path.
    """
    from pptx import Presentation  # type: ignore[import]
    from pptx.oxml.ns import qn  # type: ignore[import]

    # Resolve output path first so we can check for template/output collision
    if output_path is None:
        output_path = Path.cwd() / _safe_filename(spec.title)
    else:
        if output_path.suffix.lower() != ".pptx":
            output_path = output_path.with_suffix(".pptx")

    use_template = (
        template_path is not None
        and template_path.exists()
        and template_path.resolve() != output_path.resolve()
    )

    if use_template:
        prs = Presentation(str(template_path))
        _clear_slides(prs, qn)
    else:
        prs = Presentation()

    # Build a name→layout mapping (exact then case-insensitive)
    layout_map: dict[str, object] = {}
    for layout in prs.slide_layouts:
        layout_map[layout.name] = layout
        layout_map[layout.name.lower()] = layout

    for slide_spec in sorted(spec.slides, key=lambda s: s.index):
        _add_slide(prs, slide_spec, layout_map)

    prs.save(str(output_path))
    return output_path


def _clear_slides(prs: object, qn: object) -> None:
    """Remove all slides from a loaded Presentation, releasing their ZIP parts."""
    sldIdLst = prs.element.find(qn("p:sldIdLst"))  # type: ignore[attr-defined]
    if sldIdLst is None:
        return
    for sldId in list(sldIdLst):
        rId = sldId.get(f"{{{_R_NS}}}id")
        if rId:
            try:
                prs.part.drop_rel(rId)  # type: ignore[attr-defined]
            except Exception:
                pass
        sldIdLst.remove(sldId)


def _add_slide(
    prs: object,
    slide_spec: SlideSpec,
    layout_map: dict[str, object],
) -> None:
    # Resolve layout: exact name → lowercase → first available
    layout = (
        layout_map.get(slide_spec.layout)
        or layout_map.get(slide_spec.layout.lower())
        or next(iter(layout_map.values()), None)
    )

    slide = prs.slides.add_slide(layout)  # type: ignore[attr-defined]

    # Title placeholder (idx 0)
    if slide.shapes.title:
        slide.shapes.title.text = slide_spec.title

    # Body/content placeholder (first non-title placeholder)
    if slide_spec.bullets:
        body_ph = _find_body_placeholder(slide)
        if body_ph is not None:
            tf = body_ph.text_frame
            tf.clear()
            for i, bullet in enumerate(slide_spec.bullets):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    tf.add_paragraph().text = bullet

    # Speaker notes
    if slide_spec.notes:
        slide.notes_slide.notes_text_frame.text = slide_spec.notes


def _find_body_placeholder(slide: object) -> object | None:
    """Return the first non-title placeholder (idx != 0), or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:
            return ph
    return None
