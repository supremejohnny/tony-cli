from __future__ import annotations

import hashlib
import json
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from .mock_client import LLMClient
    from .workspace import WorkspaceContext

from .prompts_catalog import catalog_system_prompt, catalog_user_prompt
from .theme_extractor import extract_theme_tokens

# Only template PPTX files are cataloged
_PPTX_EXT = ".pptx"


# ---------------------------------------------------------------------------
# Shape formatting helpers
# ---------------------------------------------------------------------------

def _emu(v: int) -> float:
    return round(v / 914400, 2)


def _shape_tag(shape) -> str:
    """Classify a shape into TEXT / IMAGE / GROUP / DECO."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import]
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "IMAGE"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "GROUP"
    if shape.has_text_frame and shape.text_frame.text.strip():
        return "TEXT"
    return "DECO"


def _format_slides_for_prompt(path: Path) -> str:
    """Produce a compact text representation of all slides for the LLM prompt."""
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(str(path))
    lines: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f'Slide {i} | layout: "{slide.slide_layout.name}"')
        for shape in slide.shapes:
            tag = _shape_tag(shape)
            name = shape.name
            pos = (
                f"pos=({_emu(shape.left)},{_emu(shape.top)}) "
                f"{_emu(shape.width)}x{_emu(shape.height)}in"
            )
            if tag == "TEXT":
                text = shape.text_frame.text.strip()[:80].replace("\n", " | ")
                lines.append(f'  [TEXT]  "{name}"  {pos} → "{text}"')
            elif tag == "IMAGE":
                lines.append(f'  [IMAGE] "{name}"  {pos}')
            elif tag == "GROUP":
                lines.append(f'  [GROUP] "{name}"  {pos}')
            else:
                lines.append(f'  [DECO]  "{name}"  {pos}')
        lines.append("")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Programmatic shape-name extraction (for validation)
# ---------------------------------------------------------------------------

def _collect_valid_shape_names(path: Path) -> dict[int, set[str]]:
    """Return {1-based slide number → set of actual shape names} from the PPTX."""
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(str(path))
    return {
        i: {shape.name for shape in slide.shapes}
        for i, slide in enumerate(prs.slides, 1)
    }


def _validate_catalog_slots(
    patterns: list[dict],
    valid_names: dict[int, set[str]],
) -> list[dict]:
    """Drop any slot whose shape_name does not exist on the declared source_slide.

    Prints a warning for each dropped slot so the user can diagnose prompt issues.
    """
    for p in patterns:
        src = p.get("source_slide", 0)
        allowed = valid_names.get(src, set())
        good: list[dict] = []
        for slot in p.get("slots", []):
            sn = slot.get("shape_name", "")
            if sn in allowed:
                good.append(slot)
            else:
                pid = p.get("slide_id") or p.get("pattern_id") or f"slide_{src}"
                print(
                    f"  [catalog] slide {src} ({pid}): shape {sn!r} not found — slot dropped. "
                    f"Valid names: {sorted(allowed)}"
                )
        p["slots"] = good
    return patterns


# ---------------------------------------------------------------------------
# Response parsing
# ---------------------------------------------------------------------------

def _parse_catalog_response(raw: str, label: str) -> list[dict]:
    """Extract a JSON array from the LLM response with 3-step fallback."""
    text = raw.strip()

    # Step 1: direct parse
    try:
        result = json.loads(text)
        if isinstance(result, list):
            return result
    except json.JSONDecodeError:
        pass

    # Step 2: ```json [...] ``` fence
    match = re.search(r"```(?:json)?\s*(\[.*?\])\s*```", text, re.DOTALL)
    if match:
        try:
            result = json.loads(match.group(1))
            if isinstance(result, list):
                return result
        except json.JSONDecodeError:
            pass

    # Step 3: first bare [ ... ] block
    match = re.search(r"\[.*\]", text, re.DOTALL)
    if match:
        try:
            result = json.loads(match.group(0))
            if isinstance(result, list):
                return result
        except json.JSONDecodeError:
            pass

    raise ValueError(
        f"Could not parse {label} catalog response as JSON array.\n"
        f"Response preview:\n{raw[:500]}"
    )


# ---------------------------------------------------------------------------
# Hashing / cache check
# ---------------------------------------------------------------------------

def _compute_file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return f"sha256:{h.hexdigest()}"


def _is_current(catalog_path: Path, file_hash: str) -> bool:
    if not catalog_path.exists():
        return False
    try:
        existing = json.loads(catalog_path.read_text(encoding="utf-8"))
        return existing.get("source", {}).get("file_hash", "") == file_hash
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Schema normalization + deduplication
# ---------------------------------------------------------------------------

def _deduplicate_patterns(patterns: list[dict]) -> list[dict]:
    """Merge patterns whose slots share an identical set of shape_name values.

    Keeps the entry with the lowest source_slide as the canonical pattern,
    sets reusable=True, and unions fit_for / not_fit_for across the group.
    Patterns with no slots (or empty slot lists) are kept as-is.
    """
    seen: dict[frozenset, int] = {}   # shape_name_set → index in result
    result: list[dict] = []

    for p in patterns:
        shape_key = frozenset(s.get("shape_name", "") for s in p.get("slots", []) if s.get("shape_name"))
        if not shape_key:
            result.append(p)
            continue

        if shape_key in seen:
            existing = result[seen[shape_key]]
            # Union fit_for / not_fit_for (preserve order, remove exact dupes)
            for field in ("fit_for", "not_fit_for"):
                merged = list(dict.fromkeys(existing.get(field, []) + p.get(field, [])))
                existing[field] = merged
            existing["reusable"] = True
        else:
            seen[shape_key] = len(result)
            result.append(p)

    return result


def _normalize_pattern(p: dict, slide_number: int) -> dict:
    """Ensure consistent field names regardless of what the model output."""
    # Normalize slide number field: slide_number → source_slide
    if "source_slide" not in p:
        p["source_slide"] = p.pop("slide_number", slide_number)
    # Drop unknown top-level keys that add no value downstream
    p.pop("pattern_name", None)
    return p


# ---------------------------------------------------------------------------
# Catalog readers
# ---------------------------------------------------------------------------

_SUMMARY_KEYS = {"pattern_id", "source_slide", "description", "fit_for", "not_fit_for", "reusable"}
_SLOT_SUMMARY_KEYS = {"name", "content_type"}


def load_catalog_summary(catalog_path: Path) -> list[dict]:
    """Return patterns with slim slot info — used by Phase 2 planner (legacy v2 path).

    Kept for backward compat with the `generate` command and any external callers.
    New code should use load_catalog_for_planner().
    """
    data = json.loads(catalog_path.read_text(encoding="utf-8"))
    result = []
    for p in data.get("patterns", []):
        entry = {k: v for k, v in p.items() if k in _SUMMARY_KEYS}
        entry["slots"] = [
            {k: v for k, v in s.items() if k in _SLOT_SUMMARY_KEYS}
            for s in p.get("slots", [])
        ]
        result.append(entry)
    return result


def load_catalog_for_planner(catalog_path: Path) -> dict:
    """Return a slim view for the fill pipeline planner (v3).

    Returns::

        {
          "special_slides": [{"slide_id": str, "source_slide": int, "slots": [...slim...]}],
          "patterns":       [{"pattern_id": str, "source_slide": int, "reusable": bool,
                               "description": str, "fit_for": [...], "slots": [...slim...]}],
        }

    Slim slots include only ``name`` and ``content_type`` (drops shape_name / max_chars
    to keep token cost low for the planner).

    Falls back gracefully for v2 catalogs (only special_slides, no patterns).
    """
    data = json.loads(catalog_path.read_text(encoding="utf-8"))

    version = data.get("version", "2.0")
    if version == "2.0":
        print(
            "  [fill] Warning: catalog is v2.0 (no reusable patterns). "
            "Run 'powergen catalog --force' to upgrade to v3."
        )

    specials = []
    for s in data.get("special_slides", []):
        specials.append({
            "slide_id": s.get("slide_id", ""),
            "source_slide": s["source_slide"],
            "slots": [
                {"name": sl["name"], "content_type": sl["content_type"]}
                for sl in s.get("slots", [])
            ],
        })

    patterns = []
    for p in data.get("patterns", []):
        patterns.append({
            "pattern_id": p.get("pattern_id", ""),
            "source_slide": p["source_slide"],
            "reusable": p.get("reusable", True),
            "description": p.get("description", ""),
            "fit_for": p.get("fit_for", []),
            "slots": [
                {"name": sl["name"], "content_type": sl["content_type"]}
                for sl in p.get("slots", [])
            ],
        })

    return {"special_slides": specials, "patterns": patterns}


_SLOT_KEYS = {"shape_name", "name", "content_type", "max_chars"}


def load_catalog_slots(catalog_path: Path, pattern_ids: list[str]) -> list[dict]:
    """Return full slot details for *pattern_ids* only — used by Phase 3 filler (legacy).

    Only the patterns the planner selected are loaded, and slot descriptions are
    stripped (shape_name / name / content_type / max_chars are sufficient for filling).
    """
    data = json.loads(catalog_path.read_text(encoding="utf-8"))
    id_set = set(pattern_ids)
    result = []
    for p in data.get("patterns", []):
        if p.get("pattern_id") not in id_set:
            continue
        p = dict(p)
        p["slots"] = [{k: v for k, v in s.items() if k in _SLOT_KEYS} for s in p.get("slots", [])]
        result.append(p)
    return result


# ---------------------------------------------------------------------------
# v2/v3 catalog readers (used by dynamic_renderer and generate command)
# ---------------------------------------------------------------------------

def load_special_slides_meta(catalog_path: Path) -> dict[str, dict]:
    """Return {slide_id: {"source_slide": int, "slots": {name: slot_dict}}}

    Used by dynamic_renderer to locate special slides and fill their slots.
    Works for both v2 and v3 catalogs.
    """
    data = json.loads(catalog_path.read_text(encoding="utf-8"))
    result: dict[str, dict] = {}
    for s in data.get("special_slides", []):
        slide_id = s.get("slide_id", "")
        if not slide_id:
            continue
        result[slide_id] = {
            "source_slide": s["source_slide"],
            "slots": {slot["name"]: slot for slot in s.get("slots", [])},
        }
    return result


def load_catalog_theme(catalog_path: Path) -> dict:
    """Return the theme tokens dict from a v2/v3 catalog."""
    data = json.loads(catalog_path.read_text(encoding="utf-8"))
    return data.get("theme", {})


def load_available_special_slide_ids(catalog_path: Path) -> list[str]:
    """Return sorted list of slide_id values from a v2/v3 catalog."""
    data = json.loads(catalog_path.read_text(encoding="utf-8"))
    return sorted(s["slide_id"] for s in data.get("special_slides", []) if s.get("slide_id"))


# ---------------------------------------------------------------------------
# Core extraction (v3)
# ---------------------------------------------------------------------------

def extract_catalog(path: Path, client: "LLMClient") -> dict:
    """Analyze a template PPTX and return its Catalog (v3) as a dict.

    v3 schema:
      - ``theme``          — visual tokens extracted without LLM
      - ``special_slides`` — reusable=False slides (cover, profile, etc.) with slots
      - ``patterns``       — reusable=True slides (section headers, card layouts, etc.)

    Shape-name validation: all slot shape_names are cross-checked against
    programmatically extracted shape names. Invalid names are dropped with a warning.
    """
    slides_repr = _format_slides_for_prompt(path)
    raw = client.generate(catalog_system_prompt(), catalog_user_prompt(path.name, slides_repr))
    patterns = _parse_catalog_response(raw, path.name)

    # Normalize: ensure source_slide field exists
    patterns = [_normalize_pattern(p, i + 1) for i, p in enumerate(patterns)]

    # Validate shape names against real PPTX (critical fix for hallucination)
    valid_names = _collect_valid_shape_names(path)
    patterns = _validate_catalog_slots(patterns, valid_names)

    # Split into special_slides (reusable=False with slide_id) and
    # reusable patterns (reusable=True with pattern_id)
    special_slides: list[dict] = []
    reusable_patterns: list[dict] = []

    for p in patterns:
        is_reusable = p.get("reusable", True)

        if not is_reusable:
            # Special or keep slide
            entry: dict = {
                "source_slide": p["source_slide"],
                "slide_id": p.get("slide_id") or p.get("pattern_id", f"special_{p['source_slide']:02d}"),
                "slots": p.get("slots", []),
            }
            special_slides.append(entry)
        else:
            entry = {
                "pattern_id": p.get("pattern_id", f"pattern_{p['source_slide']:02d}"),
                "source_slide": p["source_slide"],
                "reusable": True,
                "description": p.get("description", ""),
                "fit_for": p.get("fit_for", []),
                "slots": p.get("slots", []),
            }
            reusable_patterns.append(entry)

    # Deduplicate reusable patterns with identical slot shapes
    reusable_patterns = _deduplicate_patterns(reusable_patterns)

    theme = extract_theme_tokens(path)

    return {
        "version": "3.0",
        "source": {
            "file_name": path.name,
            "file_hash": _compute_file_hash(path),
            "analyzed_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        },
        "theme": theme,
        "special_slides": special_slides,
        "patterns": reusable_patterns,
    }


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def run_catalog(
    workspace: "WorkspaceContext",
    client: "LLMClient",
    catalog_dir: Path,
    force: bool = False,
) -> None:
    """Catalog all template PPTX files in *workspace*, writing to *catalog_dir*."""
    catalog_dir.mkdir(parents=True, exist_ok=True)

    candidates = [ti.path for ti in workspace.templates if ti.path.suffix.lower() == _PPTX_EXT]

    if not candidates:
        print("No PPTX template files found to catalog.")
        return

    print(f"Scanning {len(candidates)} template(s)...")

    for i, path in enumerate(candidates, 1):
        catalog_path = catalog_dir / (path.stem + ".catalog.json")
        prefix = f"[{i}/{len(candidates)}] {path.name}"

        if not force:
            file_hash = _compute_file_hash(path)
            if _is_current(catalog_path, file_hash):
                print(f"{prefix} -- unchanged, skipping.")
                continue

        print(f"{prefix} -> cataloging...")
        try:
            data = extract_catalog(path, client)
        except Exception as exc:
            print(f"  Error: {exc}")
            continue

        catalog_path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        print(f"  Saved: {catalog_path.name}")

    print("Done.")
