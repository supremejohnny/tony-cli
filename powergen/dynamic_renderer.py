from __future__ import annotations

import shutil
from copy import deepcopy
from pathlib import Path

from pptx.util import Inches, Pt, Emu  # type: ignore[import]
from pptx.dml.color import RGBColor  # type: ignore[import]
from pptx.enum.text import PP_ALIGN  # type: ignore[import]
from pptx.oxml.ns import qn  # type: ignore[import]

from .catalog import load_special_slides_meta
from .catalog_filler import _fill_slide, _reorder_slides
from .theme_extractor import extract_theme_tokens


# ---------------------------------------------------------------------------
# Layout constants (from pptx skill typography guidelines)
# ---------------------------------------------------------------------------

MARGIN = Inches(0.5)        # slide edge margin
GAP    = Inches(0.3)        # gap between content blocks

FONT_TITLE_PT   = 40
FONT_SECTION_PT = 44
FONT_HEADER_PT  = 18
FONT_BODY_PT    = 15
FONT_CAPTION_PT = 11


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def _parse_rgb(hex_color: str) -> RGBColor:
    """Parse '#RRGGBB' → RGBColor."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "2E75B6"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _is_dark(hex_color: str) -> bool:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return False
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return luminance < 128


# ---------------------------------------------------------------------------
# Shape creation helpers
# ---------------------------------------------------------------------------

def _add_textbox(slide, left, top, width, height,
                 text: str, font_size: int, bold: bool,
                 font_name: str, color: RGBColor,
                 align=PP_ALIGN.LEFT, wrap: bool = True) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color


def _set_slide_background(slide, hex_color: str) -> None:
    from pptx.oxml import parse_xml  # type: ignore[import]
    from lxml import etree  # type: ignore[import]

    rgb = _parse_rgb(hex_color)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_filled_rect(slide, left, top, width, height, hex_color: str) -> None:
    from pptx.util import Emu as _Emu  # type: ignore[import]
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _parse_rgb(hex_color)
    shape.line.fill.background()  # no border


# ---------------------------------------------------------------------------
# Slide width / height helpers
# ---------------------------------------------------------------------------

def _slide_dims(prs):
    return prs.slide_width, prs.slide_height


# ---------------------------------------------------------------------------
# Find blank layout
# ---------------------------------------------------------------------------

def _find_blank_layout(prs):
    """Return the slide layout best suited as a blank content canvas.

    Prefers layouts that have *no explicit background override* so the slide
    master's visual decorations (logo, borders, background color/image) show
    through on every new slide.  Among those, picks the one with the fewest
    content placeholders.
    """
    no_bg: list[tuple[int, object]] = []
    has_bg: list[tuple[int, object]] = []

    for layout in prs.slide_layouts:
        ph_count = len(list(layout.placeholders))
        has_override = layout._element.find(qn("p:bg")) is not None
        (has_bg if has_override else no_bg).append((ph_count, layout))

    candidates = no_bg if no_bg else has_bg
    candidates.sort(key=lambda x: x[0])
    return candidates[0][1]


# ---------------------------------------------------------------------------
# Clone template slide as visual canvas
# ---------------------------------------------------------------------------

def _clone_slide_as_canvas(prs, source_idx: int):
    """Clone a template slide stripped of all content placeholders.

    Copies the source slide's background fill and every *non-placeholder*
    shape (images, coloured rectangles, border lines, logos) to a new blank
    slide.  The result is a clean visual canvas that carries the template's
    full visual identity — the caller then adds content text-boxes on top.
    """
    from copy import deepcopy

    source = list(prs.slides)[source_idx]
    blank_layout = _find_blank_layout(prs)
    new_slide = prs.slides.add_slide(blank_layout)

    # ── Replace shape tree ────────────────────────────────────────────────
    new_sp_tree = new_slide.shapes._spTree
    src_sp_tree = source.shapes._spTree

    # Remove all non-structural children from the new slide's shape tree
    for child in list(new_sp_tree):
        if child.tag not in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            new_sp_tree.remove(child)

    _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _has_text(sp_elem) -> bool:
        """Return True if this <p:sp> contains any non-whitespace text."""
        txBody = sp_elem.find(qn("p:txBody"))
        if txBody is None:
            return False
        for t_el in txBody.iter(f"{{{_NS_A}}}t"):
            if t_el.text and t_el.text.strip():
                return True
        return False

    # Copy decorative shapes from template.
    # Skip ALL <p:sp> that contain text (whether placeholder or plain text box) —
    # those are content, not decoration.  Keep images, connectors, filled
    # rectangles, group shapes, etc. which form the visual identity.
    for child in src_sp_tree:
        if child.tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            continue  # structural — already present

        if child.tag == qn("p:sp") and _has_text(child):
            continue  # content shape with real text — omit

        new_sp_tree.append(deepcopy(child))

    # ── Copy slide background fill ────────────────────────────────────────
    src_cSld = source._element.find(qn("p:cSld"))
    dst_cSld = new_slide._element.find(qn("p:cSld"))
    if src_cSld is not None and dst_cSld is not None:
        src_bg = src_cSld.find(qn("p:bg"))
        if src_bg is not None:
            dst_bg = dst_cSld.find(qn("p:bg"))
            if dst_bg is not None:
                dst_cSld.remove(dst_bg)
            dst_cSld.insert(0, deepcopy(src_bg))

    return new_slide


# ---------------------------------------------------------------------------
# Generic slide renderers
# ---------------------------------------------------------------------------

def _render_section_divider(slide, entry: dict, theme: dict, prs) -> None:
    sw, sh = _slide_dims(prs)
    bg = theme.get("bg_color", "#1B2A4A")
    acc = theme.get("accent_color", "#2E75B6")
    hfont = theme.get("heading_font", "Calibri")

    # Use bg_color for the divider background; if it's light, use accent instead
    div_bg = bg if _is_dark(bg) else acc
    _set_slide_background(slide, div_bg)

    text_color = RGBColor(0xFF, 0xFF, 0xFF)
    title = entry.get("title", "")
    subtitle = entry.get("subtitle", "")

    # Vertical center
    content_h = Pt(FONT_SECTION_PT).emu + (Pt(FONT_BODY_PT).emu + GAP if subtitle else 0)
    top = (sh - content_h) // 2

    _add_textbox(
        slide,
        MARGIN, top, sw - 2 * MARGIN, Pt(FONT_SECTION_PT).emu + GAP,
        title, FONT_SECTION_PT, True, hfont, text_color, PP_ALIGN.CENTER,
    )
    if subtitle:
        _add_textbox(
            slide,
            MARGIN, top + Pt(FONT_SECTION_PT).emu + GAP,
            sw - 2 * MARGIN, Pt(FONT_BODY_PT).emu * 2,
            subtitle, FONT_BODY_PT, False, hfont, text_color, PP_ALIGN.CENTER,
        )


def _render_content_simple(slide, entry: dict, theme: dict, prs) -> None:
    sw, sh = _slide_dims(prs)
    hfont = theme.get("heading_font", "Calibri")
    bfont = theme.get("body_font", "Calibri")
    acc = theme.get("accent_color", "#2E75B6")
    dark_text = RGBColor(0x1A, 0x1A, 0x1A)
    acc_rgb = _parse_rgb(acc)

    title = entry.get("title", "")
    bullets: list[str] = entry.get("bullets", [])

    # Title
    _add_textbox(
        slide,
        MARGIN, MARGIN, sw - 2 * MARGIN, Pt(FONT_TITLE_PT).emu + GAP,
        title, FONT_TITLE_PT, True, hfont, acc_rgb,
    )

    # Bullet list
    bullet_top = MARGIN + Pt(FONT_TITLE_PT).emu + GAP * 2
    bullet_h = sh - bullet_top - MARGIN

    txBox = slide.shapes.add_textbox(MARGIN, bullet_top, sw - 2 * MARGIN, bullet_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(FONT_BODY_PT)
        run.font.name = bfont
        run.font.color.rgb = dark_text


def _render_content_structured(slide, entry: dict, theme: dict, prs) -> None:
    sw, sh = _slide_dims(prs)
    hfont = theme.get("heading_font", "Calibri")
    bfont = theme.get("body_font", "Calibri")
    acc = theme.get("accent_color", "#2E75B6")
    acc_rgb = _parse_rgb(acc)
    dark_text = RGBColor(0x1A, 0x1A, 0x1A)
    white = RGBColor(0xFF, 0xFF, 0xFF)

    title = entry.get("title", "")
    points: list[dict] = entry.get("points", [])
    n = len(points)

    # Title
    _add_textbox(
        slide,
        MARGIN, MARGIN, sw - 2 * MARGIN, Pt(FONT_TITLE_PT).emu + GAP,
        title, FONT_TITLE_PT, True, hfont, acc_rgb,
    )

    if n == 0:
        return

    cards_top = MARGIN + Pt(FONT_TITLE_PT).emu + GAP * 2
    cards_h = sh - cards_top - MARGIN

    if n <= 2:
        # Side-by-side comparison
        card_w = (sw - 2 * MARGIN - GAP) // 2
        for i, pt in enumerate(points[:2]):
            cx = MARGIN + i * (card_w + GAP)
            _add_filled_rect(slide, cx, cards_top, card_w, Pt(FONT_HEADER_PT).emu + GAP, acc)
            _add_textbox(slide, cx + GAP // 2, cards_top + GAP // 4,
                         card_w - GAP, Pt(FONT_HEADER_PT).emu + GAP // 2,
                         pt.get("title", ""), FONT_HEADER_PT, True, hfont, white)
            desc_top = cards_top + Pt(FONT_HEADER_PT).emu + GAP
            _add_textbox(slide, cx + GAP // 2, desc_top,
                         card_w - GAP, cards_h - Pt(FONT_HEADER_PT).emu - GAP,
                         pt.get("desc", ""), FONT_BODY_PT, False, bfont, dark_text)
    else:
        # Grid (2-column)
        cols = 2
        rows = (n + 1) // 2
        card_w = (sw - 2 * MARGIN - GAP) // cols
        card_h = (cards_h - GAP * (rows - 1)) // rows
        header_h = Pt(FONT_HEADER_PT).emu + GAP

        for i, pt in enumerate(points):
            row, col = divmod(i, cols)
            cx = MARGIN + col * (card_w + GAP)
            cy = cards_top + row * (card_h + GAP)

            _add_filled_rect(slide, cx, cy, card_w, header_h, acc)
            _add_textbox(slide, cx + GAP // 2, cy + GAP // 4,
                         card_w - GAP, header_h,
                         pt.get("title", ""), FONT_HEADER_PT, True, hfont, white)
            _add_textbox(slide, cx + GAP // 2, cy + header_h + GAP // 2,
                         card_w - GAP, card_h - header_h - GAP,
                         pt.get("desc", ""), FONT_BODY_PT, False, bfont, dark_text)


def _render_two_column(slide, entry: dict, theme: dict, prs) -> None:
    sw, sh = _slide_dims(prs)
    hfont = theme.get("heading_font", "Calibri")
    bfont = theme.get("body_font", "Calibri")
    acc = theme.get("accent_color", "#2E75B6")
    acc_rgb = _parse_rgb(acc)
    dark_text = RGBColor(0x1A, 0x1A, 0x1A)

    title = entry.get("title", "")
    left_col = entry.get("left", {})
    right_col = entry.get("right", {})

    _add_textbox(
        slide,
        MARGIN, MARGIN, sw - 2 * MARGIN, Pt(FONT_TITLE_PT).emu + GAP,
        title, FONT_TITLE_PT, True, hfont, acc_rgb,
    )

    col_top = MARGIN + Pt(FONT_TITLE_PT).emu + GAP * 2
    col_h = sh - col_top - MARGIN
    col_w = (sw - 2 * MARGIN - GAP) // 2

    for i, col in enumerate([left_col, right_col]):
        cx = MARGIN + i * (col_w + GAP)
        heading = col.get("heading", "")
        bullets: list[str] = col.get("bullets", [])

        _add_textbox(slide, cx, col_top, col_w, Pt(FONT_HEADER_PT).emu + GAP,
                     heading, FONT_HEADER_PT, True, hfont, acc_rgb)

        body_top = col_top + Pt(FONT_HEADER_PT).emu + GAP
        txBox = slide.shapes.add_textbox(cx, body_top, col_w, col_h - Pt(FONT_HEADER_PT).emu - GAP)
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(3)
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.size = Pt(FONT_BODY_PT)
            run.font.name = bfont
            run.font.color.rgb = dark_text


def _render_timeline(slide, entry: dict, theme: dict, prs) -> None:
    sw, sh = _slide_dims(prs)
    hfont = theme.get("heading_font", "Calibri")
    bfont = theme.get("body_font", "Calibri")
    acc = theme.get("accent_color", "#2E75B6")
    acc_rgb = _parse_rgb(acc)
    dark_text = RGBColor(0x1A, 0x1A, 0x1A)
    white = RGBColor(0xFF, 0xFF, 0xFF)

    title = entry.get("title", "")
    steps: list[dict] = entry.get("steps", [])

    _add_textbox(
        slide,
        MARGIN, MARGIN, sw - 2 * MARGIN, Pt(FONT_TITLE_PT).emu + GAP,
        title, FONT_TITLE_PT, True, hfont, acc_rgb,
    )

    steps_top = MARGIN + Pt(FONT_TITLE_PT).emu + GAP * 2
    steps_h = sh - steps_top - MARGIN
    n = len(steps) or 1
    step_h = (steps_h - GAP * (n - 1)) // n
    num_size = min(Inches(0.45), step_h)

    for i, step in enumerate(steps):
        sy = steps_top + i * (step_h + GAP)
        # Number circle (approximated as filled square for python-pptx)
        _add_filled_rect(slide, MARGIN, sy, num_size, num_size, acc)
        _add_textbox(
            slide,
            MARGIN + Pt(2).emu, sy + Pt(2).emu, num_size - Pt(4).emu, num_size - Pt(4).emu,
            str(i + 1), FONT_HEADER_PT, True, hfont, white, PP_ALIGN.CENTER, wrap=False,
        )
        # Step label + description
        text_left = MARGIN + num_size + GAP
        text_w = sw - text_left - MARGIN
        label = step.get("label", "")
        desc = step.get("desc", "")
        _add_textbox(slide, text_left, sy, text_w, Pt(FONT_HEADER_PT).emu + GAP,
                     label, FONT_HEADER_PT, True, hfont, dark_text)
        if desc:
            _add_textbox(slide, text_left, sy + Pt(FONT_HEADER_PT).emu + GAP // 2,
                         text_w, step_h - Pt(FONT_HEADER_PT).emu - GAP,
                         desc, FONT_BODY_PT, False, bfont, dark_text)


def _render_generic_slide(slide, entry: dict, theme: dict, prs) -> None:
    slide_type = entry.get("type", "")
    if slide_type == "section_divider":
        _render_section_divider(slide, entry, theme, prs)
    elif slide_type == "content_simple":
        _render_content_simple(slide, entry, theme, prs)
    elif slide_type == "content_structured":
        points = entry.get("points", [])
        if len(points) >= 5:
            # Downgrade to bullet list
            bullets = [f'{p.get("title", "")}: {p.get("desc", "")}' for p in points]
            _render_content_simple(slide, {**entry, "bullets": bullets}, theme, prs)
        else:
            _render_content_structured(slide, entry, theme, prs)
    elif slide_type == "two_column":
        _render_two_column(slide, entry, theme, prs)
    elif slide_type == "timeline":
        _render_timeline(slide, entry, theme, prs)
    # Unknown types: leave slide blank (safe fallback)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def render_typed_plan(
    plan: list[dict],
    template_path: Path,
    catalog_path: Path,
    output_path: Path,
) -> Path:
    """Render a typed slide plan to a PPTX file.

    Special slides (title, profile, etc.) are deep-copied from the template
    and their slots filled. Generic slides are created from scratch using the
    template's visual theme tokens.

    Returns *output_path*.
    """
    from pptx import Presentation  # type: ignore[import]

    # Work on a copy so the template is never mutated
    shutil.copy2(str(template_path), str(output_path))
    prs = Presentation(str(output_path))

    # Load catalog data — extract theme fresh from template (not from catalog cache)
    theme = extract_theme_tokens(template_path)
    special_meta = load_special_slides_meta(catalog_path)

    # Build lookup: slide_id → 0-based source slide index
    special_map: dict[str, int] = {
        sid: meta["source_slide"] - 1
        for sid, meta in special_meta.items()
    }

    # Identify a representative content slide to use as visual canvas base.
    # Pick the first template slide that is NOT a special slide — it carries
    # the template's background, logo, borders, etc.
    special_indices = set(special_map.values())
    n_template_slides = len(list(prs.slides))
    content_candidates = [i for i in range(n_template_slides) if i not in special_indices]
    content_base_idx = content_candidates[0] if content_candidates else 0

    # First pass: add all new slides, track ordered indices
    ordered_indices: list[int] = []
    special_to_fill: list[tuple[int, dict, dict]] = []  # (src_idx, entry, slot_meta)

    for entry in plan:
        slide_type = entry.get("type", "")

        if slide_type in ("title", "special"):
            sid = entry.get("special_slide", "")
            src_idx = special_map.get(sid)
            if src_idx is None:
                print(f"  Warning: special_slide '{sid}' not found in catalog, skipping.")
                continue
            if src_idx in {i for i in ordered_indices if i < len(special_map)}:
                # Already scheduled — reuse the same index
                pass
            ordered_indices.append(src_idx)
            slot_meta = special_meta.get(sid, {}).get("slots", {})
            special_to_fill.append((src_idx, entry, slot_meta))

        else:
            # Clone a template content slide as visual canvas, then render content on top
            new_slide = _clone_slide_as_canvas(prs, content_base_idx)
            new_idx = len(list(prs.slides)) - 1
            ordered_indices.append(new_idx)
            _render_generic_slide(new_slide, entry, theme, prs)

    # Second pass: fill special slides
    for src_idx, entry, slot_meta in special_to_fill:
        try:
            slide = prs.slides[src_idx]
        except IndexError:
            continue
        _fill_slide(slide, entry.get("slots", {}), slot_meta)

    # Reorder / prune: keep only slides in ordered_indices, in that order
    if ordered_indices:
        _reorder_slides(prs, ordered_indices)

    prs.save(str(output_path))
    return output_path
