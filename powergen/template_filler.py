from __future__ import annotations

import difflib
import json
import re
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from .mock_client import LLMClient
    from .workspace import WorkspaceContext

from .prompts_template import (
    content_mapping_system_prompt,
    content_mapping_user_prompt,
    template_analysis_system_prompt,
    template_analysis_user_prompt,
)


# ---------------------------------------------------------------------------
# Public helpers
# ---------------------------------------------------------------------------

def pick_template(templates: list, hint: str) -> object:
    """Return the best matching TemplateInfo from *templates* based on words in *hint*.

    Uses difflib word-by-word similarity against template stems so that a user
    typing "sales templte" will still match "sales-template.pptx".
    Falls back to templates[0] when no word scores >= 0.6.
    """
    if len(templates) == 1 or not hint.strip():
        return templates[0]
    stems = [t.path.stem.lower() for t in templates]
    words = hint.lower().split()
    best_score = 0.0
    best_idx = 0
    for i, stem in enumerate(stems):
        for word in words:
            score = difflib.SequenceMatcher(None, word, stem).ratio()
            if score > best_score:
                best_score = score
                best_idx = i
    if best_score >= 0.6:
        return templates[best_idx]
    return templates[0]


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def fill_template(
    brief: str,
    template_path: Path,
    output_path: Path | None,
    client: "LLMClient",
    workspace: "WorkspaceContext",
) -> Path:
    """Fill a .pptx template with AI-generated content based on the brief.

    Uses python-pptx directly for text replacement — preserves all visual
    design (colors, fonts, images, layouts) while only changing text content.

    Returns the path to the generated .pptx file.
    """
    output_path = _resolve_output_path(template_path, output_path)

    # Step 1: extract text from template for LLM context
    markitdown_text = _run_markitdown(template_path)

    # Step 2: LLM call 1 — analyse template structure
    analysis_raw = client.generate(
        template_analysis_system_prompt(),
        template_analysis_user_prompt(markitdown_text),
    )
    analysis = _parse_json_response(analysis_raw, "template analysis")

    # Step 3: LLM call 2 — generate content mappings
    mapping_raw = client.generate(
        content_mapping_system_prompt(),
        content_mapping_user_prompt(brief, analysis, workspace),
    )
    mapping_data = _parse_json_response(mapping_raw, "content mapping")
    mappings: list[dict] = mapping_data.get("mappings", [])

    # Step 4: apply replacements only to relevant slides
    relevant = _get_relevant_slide_indices(analysis)
    _apply_mappings_pptx(template_path, output_path, mappings, relevant)

    return output_path


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _resolve_output_path(template_path: Path, output_path: Path | None) -> Path:
    if output_path is None:
        stem = template_path.stem + "-filled"
        output_path = template_path.parent / (stem + ".pptx")
    else:
        if output_path.suffix.lower() != ".pptx":
            output_path = output_path.with_suffix(".pptx")
    # Avoid overwriting the template itself
    if output_path.resolve() == template_path.resolve():
        output_path = template_path.parent / (template_path.stem + "-filled.pptx")
    return output_path


def _run_markitdown(pptx_path: Path) -> str:
    try:
        from markitdown import MarkItDown  # type: ignore[import]
        md = MarkItDown()
        result = md.convert(str(pptx_path))
        return result.text_content
    except ImportError as e:
        raise RuntimeError(
            "markitdown is required for Layer 2. Install it with: pip install 'markitdown[pptx]'"
        ) from e


def _parse_json_response(raw: str, label: str) -> dict:
    """Extract and parse JSON from LLM response, handling code fences."""
    text = raw.strip()
    fence = re.search(r"```(?:json)?\s*([\s\S]+?)```", text)
    if fence:
        text = fence.group(1).strip()
    start = text.find("{")
    end = text.rfind("}")
    if start != -1 and end != -1:
        text = text[start: end + 1]
    try:
        return json.loads(text)
    except json.JSONDecodeError as e:
        raise ValueError(
            f"Failed to parse {label} JSON: {e}\nRaw response:\n{raw[:500]}"
        ) from e


def _get_relevant_slide_indices(analysis: dict) -> set[int]:
    """Return 1-based slide indices where slide_relevant is True (defaults True if absent)."""
    relevant: set[int] = set()
    for slide in analysis.get("slides", []):
        if slide.get("slide_relevant", True):
            idx = slide.get("slide_index")
            if idx is not None:
                relevant.add(int(idx))
    return relevant


def _apply_mappings_pptx(
    template_path: Path,
    output_path: Path,
    mappings: list[dict],
    relevant_slides: set[int] | None = None,
) -> None:
    """Open template with python-pptx, apply text replacements, save to output."""
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(str(template_path))

    # Group mappings by slide_index (1-based)
    by_slide: dict[int, list[dict]] = {}
    for m in mappings:
        idx = m.get("slide_index")
        if idx is not None:
            by_slide.setdefault(int(idx), []).append(m)

    for slide_idx_1based, slide_mappings in by_slide.items():
        if relevant_slides is not None and slide_idx_1based not in relevant_slides:
            print(f"  Skipping slide {slide_idx_1based} (not relevant).")
            continue
        slide_idx = slide_idx_1based - 1  # python-pptx is 0-indexed
        if slide_idx >= len(prs.slides):
            print(f"  Warning: slide {slide_idx_1based} not in template (only {len(prs.slides)} slides), skipping.")
            continue
        slide = prs.slides[slide_idx]
        _replace_in_slide_pptx(slide, slide_mappings, slide_idx_1based)

    prs.save(str(output_path))


def _replace_in_slide_pptx(slide: object, mappings: list[dict], slide_num: int) -> None:
    """Replace paragraph text in a slide, preserving run formatting."""
    # Build lookup: normalised original text → replacement
    lookup: dict[str, str] = {}
    for m in mappings:
        original = m.get("original_text", "")
        replacement = m.get("replacement_text", "")
        if original:
            lookup[_normalise(original)] = replacement

    for shape in slide.shapes:  # type: ignore[attr-defined]
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = _normalise(para.text)
            if full_text not in lookup:
                continue
            replacement = lookup[full_text]
            runs = para.runs
            if not runs:
                continue
            # Put replacement in first run, blank the rest (preserve formatting)
            runs[0].text = replacement
            for run in runs[1:]:
                run.text = ""
            del lookup[full_text]  # only replace first match per slide

    # Warn about any mappings that weren't matched
    for m in mappings:
        original = m.get("original_text", "")
        if original and _normalise(original) in lookup:
            preview = original[:60].replace("\n", "\\n")
            print(f"  Warning: text not found in slide {slide_num}: {preview!r}")


def _normalise(text: str) -> str:
    """Normalise whitespace for fuzzy paragraph matching."""
    return re.sub(r"\s+", " ", text).strip()
