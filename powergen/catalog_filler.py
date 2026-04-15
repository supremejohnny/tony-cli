from __future__ import annotations

import json
import shutil
from copy import deepcopy
from pathlib import Path

from pptx.oxml.ns import qn  # type: ignore[import]


# ---------------------------------------------------------------------------
# Shape-level text writing
# ---------------------------------------------------------------------------

def _write_shape(shape, content: str, content_type: str) -> None:
    """Write *content* into *shape*'s text frame, preserving paragraph/run formatting.

    Strategy:
    - Save the first paragraph element and first run element as format templates.
    - Remove all existing <a:p> elements from the txBody.
    - Re-insert one <a:p> per output line, each containing one <a:r> cloned from
      the saved run template (keeps font, size, colour, bold, etc.).
    """
    if not shape.has_text_frame:
        return

    lines: list[str]
    if content_type == "bullets":
        lines = [l.strip() for l in content.split("\n") if l.strip()]
    else:
        lines = [content.strip()]

    if not lines:
        return

    txBody = shape.text_frame._txBody

    # --- save templates BEFORE touching the tree ---
    existing_paras = txBody.findall(qn("a:p"))
    tmpl_p = deepcopy(existing_paras[0]) if existing_paras else None
    tmpl_r: object | None = None
    if tmpl_p is not None:
        tmpl_rs = tmpl_p.findall(qn("a:r"))
        tmpl_r = deepcopy(tmpl_rs[0]) if tmpl_rs else None

    # --- remove all existing paragraphs ---
    for p in existing_paras:
        txBody.remove(p)

    # --- insert new paragraphs ---
    from lxml import etree  # type: ignore[import]

    for line in lines:
        # Build paragraph from template or scratch
        if tmpl_p is not None:
            new_p = deepcopy(tmpl_p)
            for r in new_p.findall(qn("a:r")):
                new_p.remove(r)
        else:
            new_p = etree.Element(qn("a:p"))

        # Build run from template or scratch
        if tmpl_r is not None:
            new_r = deepcopy(tmpl_r)
            t = new_r.find(qn("a:t"))
            if t is None:
                t = etree.SubElement(new_r, qn("a:t"))
        else:
            new_r = etree.Element(qn("a:r"))
            t = etree.SubElement(new_r, qn("a:t"))

        t.text = line
        new_p.append(new_r)
        txBody.append(new_p)


def _fill_slide(slide, slot_values: dict[str, str], slot_meta: dict[str, dict]) -> None:
    """Write all slot values into the matching shapes on *slide*.

    *slot_meta* maps slot_name → {shape_name, content_type, ...}.
    Shapes are matched by their exact ``shape.name`` attribute.

    When multiple shapes share the same shape_name (e.g. two 'TextBox 15' elements),
    slots that map to the same shape_name are filled in encounter order: the first
    such slot fills the first matching shape, the second fills the second, etc.
    """
    # Build targets: shape_name → list of (content, content_type) in slot order
    targets: dict[str, list[tuple[str, str]]] = {}
    for slot_name, content in slot_values.items():
        meta = slot_meta.get(slot_name)
        if meta is None:
            continue
        shape_name = meta.get("shape_name") or meta.get("name", "")
        content_type = meta.get("content_type", "text")
        if shape_name:
            targets.setdefault(shape_name, []).append((content, content_type))

    # Track how many times each shape_name has been written to
    fill_count: dict[str, int] = {}
    filled_slots: set[str] = set()

    for shape in slide.shapes:
        sn = shape.name
        if sn not in targets:
            continue
        idx = fill_count.get(sn, 0)
        if idx < len(targets[sn]):
            content, content_type = targets[sn][idx]
            _write_shape(shape, content, content_type)
            fill_count[sn] = idx + 1
            filled_slots.add(sn)

    # Warn about any slots whose shape was never found
    unfilled = set(targets.keys()) - filled_slots
    if unfilled:
        for name in sorted(unfilled):
            print(f"  Warning: shape not found on slide: {name!r}")


# ---------------------------------------------------------------------------
# Slide ordering helpers  (python-pptx XML level)
# ---------------------------------------------------------------------------

def _reorder_slides(prs, ordered_indices: list[int]) -> None:
    """Reorder and prune slides so only *ordered_indices* (0-based) remain,
    in that order.  Slides not listed are dropped from the presentation.
    """
    sldIdLst = prs.slides._sldIdLst
    all_sldId = list(sldIdLst)          # snapshot before mutation
    n_original = len(all_sldId)

    # Remove all slide id entries
    for sid in all_sldId:
        sldIdLst.remove(sid)

    # Re-add only the ones we want, in plan order
    keep_set = set(ordered_indices)
    for idx in ordered_indices:
        sldIdLst.append(all_sldId[idx])

    # Drop package relationships for removed slides so pptx stays valid
    for idx in range(n_original):
        if idx not in keep_set:
            rId = all_sldId[idx].get(qn("r:id"))
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass


# ---------------------------------------------------------------------------
# Slide cloning helpers
# ---------------------------------------------------------------------------

def _copy_slide_rels(source_slide, new_slide) -> None:
    """Copy non-structural media/hyperlink relationships from *source_slide* to
    *new_slide* and remap r:embed / r:link / r:id attributes in the new slide's
    XML so images and other media resolve correctly.

    Structural relationship types (slideLayout, notesSlide) are skipped because
    *new_slide* already has its own layout relationship from add_slide().
    """
    _STRUCTURAL_REL_TYPES = {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
    }

    rId_map: dict[str, str] = {}

    for rel in source_slide.part.rels.values():
        if rel.reltype in _STRUCTURAL_REL_TYPES:
            continue
        try:
            if rel.is_external:
                new_rId = new_slide.part.relate_to(rel._target, rel.reltype, is_external=True)
            else:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rel.rId] = new_rId
        except Exception:
            continue

    if not rId_map:
        return

    # Patch all relationship-ID attributes in the new slide's XML tree
    _NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    r_embed = f"{{{_NS_R}}}embed"
    r_link  = f"{{{_NS_R}}}link"
    r_id    = f"{{{_NS_R}}}id"

    for elem in new_slide._element.iter():
        for attr in (r_embed, r_link, r_id):
            old = elem.get(attr)
            if old and old in rId_map:
                elem.set(attr, rId_map[old])


def _clone_slide_full(prs, source_idx: int):
    """Deep-copy the template slide at *source_idx* (0-based) into a new slide.

    Unlike ``_clone_slide_as_canvas`` (which strips text shapes), this clone
    preserves ALL shapes including text — the full visual and content state of
    the source slide is replicated.  Callers then overwrite specific text shapes
    via ``_fill_slide()``.

    Returns the new slide object (already appended to prs.slides).
    """
    source_slide = list(prs.slides)[source_idx]
    source_layout = source_slide.slide_layout

    # Add new slide with same layout (inherits master decorations automatically)
    new_slide = prs.slides.add_slide(source_layout)

    # ── Replace shape tree ────────────────────────────────────────────────
    new_sp_tree = new_slide.shapes._spTree
    src_sp_tree = source_slide.shapes._spTree

    # Remove all non-structural children from the new slide's shape tree
    for child in list(new_sp_tree):
        if child.tag not in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            new_sp_tree.remove(child)

    # Deep-copy ALL shapes from source (text + deco + images — full clone)
    for child in src_sp_tree:
        if child.tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            continue  # structural — already present
        new_sp_tree.append(deepcopy(child))

    # ── Copy slide-level background ───────────────────────────────────────
    src_cSld = source_slide._element.find(qn("p:cSld"))
    dst_cSld = new_slide._element.find(qn("p:cSld"))
    if src_cSld is not None and dst_cSld is not None:
        src_bg = src_cSld.find(qn("p:bg"))
        if src_bg is not None:
            dst_bg = dst_cSld.find(qn("p:bg"))
            if dst_bg is not None:
                dst_cSld.remove(dst_bg)
            dst_cSld.insert(0, deepcopy(src_bg))

    # ── Remap image / media relationships ────────────────────────────────
    _copy_slide_rels(source_slide, new_slide)

    return new_slide


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def fill_from_plan(
    plan: list[dict],
    template_path: Path,
    catalog_path: Path,
    output_path: Path,
) -> Path:
    """Phase 3 v3: Apply an op-based content plan to a template PPTX.

    Supported ops:

      fill_special  — fill a one-time special slide in-place
      keep          — include a template slide unchanged
      clone_pattern — deep-copy a pattern slide and fill with new content

    *plan* is the list returned by ``run_catalog_plan()``::

        [{"op": "fill_special", "slide_id": str, "slots": {...}}, ...]
        [{"op": "keep",         "source_slide": int}, ...]
        [{"op": "clone_pattern","pattern_id": str,  "slots": {...}}, ...]

    Returns *output_path*.
    """
    from pptx import Presentation  # type: ignore[import]

    # ── Load catalog ─────────────────────────────────────────────────────
    catalog_data = json.loads(catalog_path.read_text(encoding="utf-8"))

    # Build lookup dicts
    special_by_id: dict[str, dict] = {
        s["slide_id"]: s for s in catalog_data.get("special_slides", [])
    }
    pattern_by_id: dict[str, dict] = {
        p["pattern_id"]: p for p in catalog_data.get("patterns", [])
    }

    # Build slot-meta lookups: {id → {slot_name → slot_dict}}
    special_slot_meta: dict[str, dict[str, dict]] = {
        sid: {sl["name"]: sl for sl in s.get("slots", [])}
        for sid, s in special_by_id.items()
    }
    pattern_slot_meta: dict[str, dict[str, dict]] = {
        pid: {sl["name"]: sl for sl in p.get("slots", [])}
        for pid, p in pattern_by_id.items()
    }

    # ── Work on a copy of the template ───────────────────────────────────
    shutil.copy2(str(template_path), str(output_path))
    prs = Presentation(str(output_path))
    n_slides = len(prs.slides)

    ordered_indices: list[int] = []   # final slide order (0-based)

    for entry in plan:
        op = entry.get("op", "")

        # ── fill_special ─────────────────────────────────────────────────
        if op == "fill_special":
            sid = entry.get("slide_id", "")
            special_meta = special_by_id.get(sid)
            if special_meta is None:
                print(f"  Warning: special slide '{sid}' not in catalog — skipping.")
                continue
            src_idx = special_meta["source_slide"] - 1
            if src_idx < 0 or src_idx >= n_slides:
                print(f"  Warning: fill_special source_slide {src_idx + 1} out of range — skipping.")
                continue
            slide = prs.slides[src_idx]
            _fill_slide(slide, entry.get("slots", {}), special_slot_meta.get(sid, {}))
            ordered_indices.append(src_idx)

        # ── keep ─────────────────────────────────────────────────────────
        elif op == "keep":
            src_slide = entry.get("source_slide", 0)
            src_idx = src_slide - 1
            if src_idx < 0 or src_idx >= n_slides:
                print(f"  Warning: keep source_slide {src_slide} out of range — skipping.")
                continue
            ordered_indices.append(src_idx)

        # ── clone_pattern ─────────────────────────────────────────────────
        elif op == "clone_pattern":
            pid = entry.get("pattern_id", "")
            pattern_meta = pattern_by_id.get(pid)
            if pattern_meta is None:
                print(f"  Warning: pattern_id '{pid}' not in catalog — skipping.")
                continue
            src_idx = pattern_meta["source_slide"] - 1
            if src_idx < 0 or src_idx >= n_slides:
                print(f"  Warning: clone_pattern source_slide {src_idx + 1} out of range — skipping.")
                continue
            new_slide = _clone_slide_full(prs, src_idx)
            new_idx = len(list(prs.slides)) - 1
            _fill_slide(new_slide, entry.get("slots", {}), pattern_slot_meta.get(pid, {}))
            ordered_indices.append(new_idx)

        else:
            print(f"  Warning: unknown op '{op}' — skipping.")

    # ── Reorder and prune ─────────────────────────────────────────────────
    if ordered_indices:
        _reorder_slides(prs, ordered_indices)

    prs.save(str(output_path))
    return output_path
