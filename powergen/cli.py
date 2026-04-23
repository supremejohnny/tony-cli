from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        prog="powergen",
        description="PowerGen — AI presentation generator",
    )
    parser.add_argument(
        "--mock",
        action="store_true",
        help="Use mock LLM client (no API calls, for testing)",
    )
    parser.add_argument(
        "--model",
        default="claude-sonnet-4-6",
        metavar="MODEL",
        help="Anthropic model to use (default: claude-sonnet-4-6)",
    )

    sub = parser.add_subparsers(dest="cmd")

    create_p = sub.add_parser("create", help="Generate a presentation plan")
    create_p.add_argument("topic", help="Topic or description of the presentation")

    revise_p = sub.add_parser("revise", help="Revise the current plan with feedback")
    revise_p.add_argument("feedback", help="What to change in the plan")

    sub.add_parser("approve", help="Approve plan and build slide spec")

    render_p = sub.add_parser("render", help="Render the approved spec to .pptx")
    render_p.add_argument(
        "--output", "-o",
        default=None,
        metavar="FILE",
        help="Output .pptx path (default: <title>.pptx in current directory)",
    )

    tmpl_p = sub.add_parser(
        "template",
        help="Compose a presentation from a pptx template (Layer 2)",
    )
    tmpl_src = tmpl_p.add_mutually_exclusive_group()
    tmpl_src.add_argument("--pptx", default=None, metavar="FILE",
                          help="Source template .pptx (schema auto-generated + cached alongside)")
    tmpl_src.add_argument("--schema", default=None, metavar="SCHEMA",
                          help="Path to a hand-authored .schema.json (advanced)")
    tmpl_p.add_argument("--topic", default=None, metavar="TOPIC",
                        help="Topic / content description for the Composer LLM")
    tmpl_p.add_argument("--plan", default=None, metavar="FILE",
                        help="Load a pre-authored plan.json instead of calling the LLM")
    tmpl_p.add_argument("--output", "-o", default=None, metavar="FILE",
                        help="Output .pptx path (default: <title>.pptx)")

    sub.add_parser("status", help="Show current project stage")
    sub.add_parser("reset", help="Reset project state to INIT")

    args = parser.parse_args(argv)

    from .mock_client import make_llm_client
    from .state import ProjectState, StateError
    from .workspace import scan_workspace

    # template command is stateless — skip ProjectState/workspace loading
    if args.cmd == "template":
        client = make_llm_client(mock=args.mock, model=args.model)
        return _run_template(args, client)

    state = ProjectState.load(Path.cwd())
    workspace = scan_workspace()
    client = make_llm_client(mock=args.mock, model=args.model)

    try:
        if args.cmd == "create":
            from .planner import generate_plan
            print("Generating plan…")
            plan = generate_plan(args.topic, workspace, client, state)
            print(f"\nPlan ready — {len(plan.slide_summaries)} slides")
            print(f"Overview: {plan.overview}")
            if plan.open_questions:
                print("Open questions:")
                for q in plan.open_questions:
                    print(f"  · {q}")

        elif args.cmd == "revise":
            from .planner import revise_plan
            print("Revising plan…")
            plan = revise_plan(args.feedback, workspace, client, state)
            print(f"\nRevised — now {len(plan.slide_summaries)} slides")

        elif args.cmd == "approve":
            from .spec_builder import build_spec
            if state.plan is None:
                print("Error: no plan yet. Run 'powergen create <topic>' first.", file=sys.stderr)
                return 1
            print("[1/2] Locking plan…")
            print("[2/2] Building slide spec…")
            spec = build_spec(state.plan, workspace, client)
            state.advance_to_approved(spec)
            print(f"\nApproved — {len(spec.slides)} slides ready to render")
            print(f"Title: {spec.title}")

        elif args.cmd == "render":
            from .renderer import render_presentation
            if state.spec is None:
                print("Error: no approved spec. Run 'powergen approve' first.", file=sys.stderr)
                return 1
            template_path = None
            if state.spec.theme_reference:
                candidate = Path.cwd() / state.spec.theme_reference
                if candidate.exists():
                    template_path = candidate
            output_path = Path(args.output) if args.output else None
            out = render_presentation(state.spec, output_path=output_path, template_path=template_path)
            state.advance_to_rendered(str(out))
            print(f"\nRendered: {out}")

        elif args.cmd == "status":
            print(f"Stage: {state.stage.value}")
            if state.plan:
                print(f"Slides planned: {len(state.plan.slide_summaries)}")
            if state.spec:
                print(f"Slides in spec: {len(state.spec.slides)}")
                print(f"Title: {state.spec.title}")
            if state.output_path:
                print(f"Output: {state.output_path}")

        elif args.cmd == "reset":
            state.reset()
            print("Project reset to INIT.")

        else:
            # No subcommand → interactive REPL
            from .repl import run_repl
            run_repl(state, client, workspace)

    except StateError as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1
    except KeyboardInterrupt:
        print("\nInterrupted.")
        return 0

    return 0


def _run_template(args, client) -> int:
    from pptx import Presentation

    from .layer2.composer.planner import generate_plan, mock_plan
    from .layer2.composer.composer import compose

    # Resolve schema + src_pptx
    if args.pptx:
        from .layer2.composer.schema_gen import load_or_generate
        pptx_path = Path(args.pptx)
        if not pptx_path.exists():
            print(f"Error: pptx not found: {pptx_path}", file=sys.stderr)
            return 1
        schema, _ = load_or_generate(pptx_path)
        src_pptx = pptx_path.resolve()
    elif args.schema:
        from .layer2.composer.schema_loader import load as load_schema
        schema_path = Path(args.schema)
        if not schema_path.exists():
            print(f"Error: schema not found: {schema_path}", file=sys.stderr)
            return 1
        schema = load_schema(schema_path)
        src_pptx = schema_path.parent / schema["source_pptx"]
        if not src_pptx.exists():
            print(f"Error: source pptx not found: {src_pptx}", file=sys.stderr)
            return 1
    else:
        print("Error: provide --pptx <template.pptx> or --schema <file.schema.json>", file=sys.stderr)
        return 1

    # Determine plan
    if args.plan:
        plan_path = Path(args.plan)
        if not plan_path.exists():
            print(f"Error: plan file not found: {plan_path}", file=sys.stderr)
            return 1
        plan = json.loads(plan_path.read_text(encoding="utf-8"))
        print(f"Loaded plan: {plan_path}")
    elif args.topic:
        from .mock_client import MockLLMClient
        if isinstance(client, MockLLMClient):
            print("Mock mode: using built-in plan (--topic ignored in mock)")
            plan = mock_plan(schema)
        else:
            print("Composing plan…")
            plan = generate_plan(schema, args.topic, client)
    else:
        print("Error: provide --topic or --plan", file=sys.stderr)
        return 1

    title = plan.get("title", "presentation")
    n_slides = len(plan.get("slides", []))
    print(f"Plan: {n_slides} slides — {title!r}")

    src_prs = Presentation(str(src_pptx))
    print("Composing…")
    dest_prs = compose(schema, src_prs, plan)

    # Output path
    if args.output:
        out_path = Path(args.output)
    else:
        safe = "".join(c if c.isalnum() or c in " -_" else "_" for c in title)
        out_path = Path(safe.strip() + ".pptx")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    dest_prs.save(str(out_path))
    print(f"Saved: {out_path}  ({len(list(dest_prs.slides))} slides)")
    return 0
