from __future__ import annotations

import copy
import re
from io import BytesIO
from pathlib import Path

_INDEXED_RE = re.compile(r"^(.+)\[(\d+)\]$")


def compose(template_path: Path, plan: dict, output_path: Path) -> Path:
    """Clone slides from template per plan, fill text_map, save to output_path."""
    from pptx import Presentation  # type: ignore[import]

    src_prs = Presentation(str(template_path))

    # dest starts as a full copy of src → all image/layout/master parts already present
    buf = BytesIO()
    src_prs.save(buf)
    buf.seek(0)
    dest_prs = Presentation(buf)

    # Collect wanted slides BEFORE detaching anything from dest_prs.
    # first_use tracks which original dest slides have been allocated.
    # Duplicates (clone_again) are created via _duplicate within dest_prs
    # so their parts share the same package parts — no cross-package pollution.
    first_use: dict[int, object] = {}
    wanted: list[tuple] = []  # (slide_obj, text_map)

    n_src = len(dest_prs.slides)

    for entry in plan.get("slides", []):
        if entry.get("type") == "generated":
            from .renderers import render_generated
            slide = render_generated(dest_prs, entry)
            if slide is not None:
                wanted.append((slide, {}))
            continue

        idx = entry.get("source_slide_index", 0)
        if idx >= n_src:
            print(f"  Warning: source_slide_index {idx} out of range ({n_src} slides), skipping.")
            continue

        text_map: dict[str, str] = entry.get("text_map", {})

        if idx not in first_use:
            slide = dest_prs.slides[idx]
            first_use[idx] = slide
        else:
            slide = _duplicate_within(dest_prs, idx)

        wanted.append((slide, text_map))

    # Rebuild dest_prs slide list from wanted (detach all, reattach in order)
    _detach_all(dest_prs)
    for slide, text_map in wanted:
        _attach(dest_prs, slide)
        _fill_slide(slide, text_map)

    filled = len(wanted)
    dest_prs.save(str(output_path))
    print(f"  Composed {filled} slides")
    return output_path


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _detach_all(prs) -> None:
    """Remove all slides from the presentation's slide list (parts stay in package)."""
    from pptx.oxml.ns import qn  # type: ignore[import]

    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(qn("r:id"))
        sld_id_lst.remove(sld_id)
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:
                pass


def _attach(prs, slide) -> None:
    """Add an existing slide part back into the presentation's slide list."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT  # type: ignore[import]
    from pptx.oxml.ns import qn  # type: ignore[import]
    from lxml import etree  # type: ignore[import]

    r_id = prs.part.relate_to(slide.part, RT.SLIDE)
    sld_id_lst = prs.slides._sldIdLst
    existing = [int(el.get("id", 0)) for el in sld_id_lst]
    next_id = max(existing, default=255) + 1
    sld_id_el = etree.SubElement(sld_id_lst, qn("p:sldId"))
    sld_id_el.set("id", str(next_id))
    sld_id_el.set(qn("r:id"), r_id)


def _find_layout(prs, layout_name: str):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    return prs.slide_layouts[0]


def _duplicate_within(prs, src_idx: int):
    """Clone prs.slides[src_idx] within the same presentation (no cross-package issues)."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT  # type: ignore[import]

    src_slide = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(_find_layout(prs, src_slide.slide_layout.name))

    # Replace shape tree with source shapes
    sp_tree = new_slide.shapes._spTree
    for el in list(sp_tree):
        sp_tree.remove(el)
    for el in src_slide.shapes._spTree:
        sp_tree.append(copy.deepcopy(el))

    # Copy non-layout relationships; both slides are in the same prs so
    # relate_to() on the same Part object is idempotent — no package duplication.
    rId_map: dict[str, str] = {}
    for rId, rel in src_slide.part.rels.items():
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue
        try:
            if rel.is_external:
                new_rId = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            else:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            if new_rId != rId:
                rId_map[rId] = new_rId
        except Exception:
            pass

    if rId_map:
        _remap_rids(new_slide._element, rId_map)

    return new_slide


def _remap_rids(element, rId_map: dict[str, str]) -> None:
    from pptx.oxml.ns import qn  # type: ignore[import]

    attrs = (qn("r:embed"), qn("r:id"), qn("r:link"))
    for node in element.iter():
        for attr in attrs:
            val = node.get(attr)
            if val in rId_map:
                node.set(attr, rId_map[val])


def _resolve_name(name: str) -> tuple[str, int]:
    """Parse 'ShapeName[N]' → ('ShapeName', N). Plain name → (name, 0)."""
    m = _INDEXED_RE.match(name)
    if m:
        return m.group(1), int(m.group(2))
    return name, 0


def _fill_slide(slide, text_map: dict[str, str]) -> None:
    """Best-effort fill: supports plain names and indexed names like 'TextBox 15[1]'."""
    groups: dict[str, list] = {}
    for shape in slide.shapes:
        if shape.has_text_frame:
            groups.setdefault(shape.name, []).append(shape)

    for key, new_text in text_map.items():
        base, idx = _resolve_name(key)
        group = groups.get(base, [])
        if idx < len(group):
            _replace_text(group[idx], new_text)
        else:
            print(f"    Warning: shape '{key}' not found, skipping.")


def _replace_text(shape, new_text: str) -> None:
    """Replace all text in a shape's text frame, preserving first-run formatting."""
    tf = shape.text_frame
    if not tf.paragraphs:
        return

    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = new_text
        for run in first_para.runs[1:]:
            run.text = ""
    else:
        first_para.text = new_text

    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""
