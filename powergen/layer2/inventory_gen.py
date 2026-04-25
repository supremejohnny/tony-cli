from __future__ import annotations

from pathlib import Path


def generate(pptx_path: Path) -> dict:
    """Extract slide inventory from a .pptx file. Pure code, zero LLM tokens.

    Duplicate shape names within a slide are suffixed with [0], [1], … so the
    LLM can target each occurrence individually via text_map keys.
    """
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(str(pptx_path))
    slides = []
    for i, slide in enumerate(prs.slides):
        raw: list[dict] = []
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                preview_cells = []
                for row in list(tbl.rows)[:2]:
                    for cell in list(row.cells)[:3]:
                        t = cell.text_frame.text.strip()
                        if t:
                            preview_cells.append(t[:30])
                if preview_cells:
                    raw.append({
                        "name": shape.name,
                        "type": "table",
                        "rows": len(tbl.rows),
                        "cols": len(tbl.columns),
                        "preview": " | ".join(preview_cells[:4]),
                    })
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                raw.append({"name": shape.name, "text": text[:120]})

        # Count occurrences per name to detect duplicates (text shapes only)
        counts: dict[str, int] = {}
        for s in raw:
            if s.get("type") != "table":
                counts[s["name"]] = counts.get(s["name"], 0) + 1

        seen: dict[str, int] = {}
        shapes: list[dict] = []
        for s in raw:
            if s.get("type") == "table":
                shapes.append(s)
                continue
            base = s["name"]
            if counts[base] > 1:
                idx = seen.get(base, 0)
                seen[base] = idx + 1
                display = f"{base}[{idx}]"
            else:
                display = base
            shapes.append({"name": display, "text": s["text"]})

        slides.append({
            "index": i,
            "layout": slide.slide_layout.name,
            "shapes": shapes,
        })
    return {"slides": slides}


def format_for_prompt(inventory: dict) -> str:
    """Render inventory as a compact text block suitable for an LLM prompt."""
    lines: list[str] = []
    for slide in inventory["slides"]:
        lines.append(f"\n--- Slide {slide['index']} | layout: {slide['layout']!r} ---")
        if not slide["shapes"]:
            lines.append("  (no text shapes)")
            continue
        for shape in slide["shapes"]:
            if shape.get("type") == "table":
                lines.append(
                    f'  [TABLE] "{shape["name"]}": {shape["rows"]}×{shape["cols"]}, preview: "{shape["preview"]}"'
                )
            else:
                preview = shape["text"].replace("\n", " | ")
                lines.append(f'  "{shape["name"]}": "{preview}"')
    return "\n".join(lines)
