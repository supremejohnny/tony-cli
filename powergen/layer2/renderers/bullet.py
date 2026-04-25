from __future__ import annotations

from pptx.util import Inches, Pt  # type: ignore[import]
from pptx.dml.color import RGBColor  # type: ignore[import]
from pptx.enum.text import PP_ALIGN  # type: ignore[import]


def render_bullet(prs, entry: dict):
    """Render a simple title + bullet list slide onto a blank layout."""
    blank_layout = _find_blank_layout(prs)
    slide = prs.slides.add_slide(blank_layout)

    title_text = entry.get("title", "")
    items: list[str] = entry.get("items", [])

    W = prs.slide_width
    H = prs.slide_height

    # Title box
    if title_text:
        txb = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), W - Inches(1), Inches(1.0)
        )
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.runs[0].font.size = Pt(28)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = RGBColor(0x1F, 0x29, 0x37)

    # Bullet box
    if items:
        txb = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), W - Inches(1), H - Inches(2.0)
        )
        tf = txb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.runs[0].font.size = Pt(18)
            p.runs[0].font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            p.space_before = Pt(6)

    return slide


def _find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower() or layout.name == "空白":
            return layout
    return prs.slide_layouts[-1]
