"""Auto-generate a schema dict from a pptx by local shape extraction.

Schema version 4 additions vs v3:
- intent_tags per slide: inferred from pptx layout names (e.g. ["cover"], ["list","highlight"])
- inferred_label per slot: heuristic title/subtitle/body from position+font size
- removed compose_hints (no longer needed with two-phase pipeline)
"""
from __future__ import annotations

import json
from pathlib import Path

_THEME_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

# Template-agnostic renderer definitions; keys are structure_type values the LLM outputs.
_GENERATED_SLIDES = {
    "list": {
        "renderer": "renderers.title_bullets",
        "max_items": 8,
        "max_chars_per_item": 100,
        "notes": "Title + bullet list. Use for independent points with no strong visual grouping.",
    },
    "comparison": {
        "renderer": "renderers.two_column",
        "max_items": 2,
        "max_chars_per_item": 200,
        "notes": "Two-column layout. Use when comparing exactly two options or perspectives.",
    },
    "cards": {
        "renderer": "renderers.card_grid",
        "max_items": 4,
        "max_chars_per_item": 80,
        "notes": "2x2 or 3x1 card grid. Use for parallel categories with equal visual weight.",
    },
    "process": {
        "renderer": "renderers.flow",
        "max_items": 5,
        "max_chars_per_item": 60,
        "notes": "Horizontal A→B→C chain. Use when items have sequential or causal dependency.",
    },
    "section": {
        "renderer": "renderers.section_divider",
        "max_items": 1,
        "max_chars_per_item": 40,
        "notes": "Section divider with large title and optional subtitle. No content items.",
    },
}

_LAYOUT_INTENT_MAP = [
    ({"封面", "cover", "title slide", "首页", "front page"}, ["cover"]),
    ({"节标题", "section header", "章节", "section divider", "节", "chapter"}, ["section"]),
    ({"目录", "agenda", "contents", "table of content", "toc", "overview"}, ["section", "list"]),
    ({"two content", "比较", "comparison", "对比", "两栏", "versus"}, ["comparison"]),
    ({"timeline", "时间线", "流程", "process", "步骤", "sequence"}, ["process"]),
    ({"结尾", "end slide", "closing", "thank", "谢谢", "结束", "goodbye"}, ["closing"]),
    ({"内容", "content", "title and content", "text and content", "bullet"}, ["list", "highlight"]),
]


def _infer_intent_tags(layout_name: str) -> list[str]:
    low = layout_name.lower()
    for keywords, tags in _LAYOUT_INTENT_MAP:
        if any(kw in low for kw in keywords):
            return list(tags)
    return []


def _infer_slot_labels(slots: dict) -> None:
    """Add inferred_label in-place: largest/topmost font → title, second → subtitle, rest → body."""
    candidates = [
        (k, v) for k, v in slots.items()
        if not v.get("visual_only") and v.get("kind") in ("text", "multiline")
    ]
    if not candidates:
        return
    sorted_c = sorted(
        candidates,
        key=lambda kv: (-(kv[1].get("font_size_pt") or 0), kv[1].get("top_pct", 0.5)),
    )
    seen_title = False
    seen_subtitle = False
    for k, v in sorted_c:
        sz = v.get("font_size_pt") or 0
        if not seen_title:
            v["inferred_label"] = "title"
            seen_title = True
        elif not seen_subtitle and sz >= 12:
            v["inferred_label"] = "subtitle"
            seen_subtitle = True
        else:
            v["inferred_label"] = "body"


_DECORATIVE_FONTS = frozenset({
    "wingdings", "wingdings 2", "wingdings 3", "webdings", "symbol",
    "marlett", "segoe ui symbol",
})


def _is_visual_only(font_name: str | None, text: str) -> bool:
    if font_name and font_name.lower() in _DECORATIVE_FONTS:
        return True
    if text:
        latin_ext = sum(1 for c in text if "Ā" <= c <= "ɏ")
        if latin_ext / max(len(text), 1) > 0.15:
            return True
    return False


def _get_font_info(shape) -> tuple[str | None, int | None]:
    """Return (font_name, font_size_pt) from first run of first paragraph."""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                sz = run.font.size
                name = run.font.name
                return name, (int(sz / 12700) if sz else None)
    except Exception:
        pass
    return None, None


def _layout_meta(shape, slide_w: int, slide_h: int) -> dict:
    try:
        return {
            "top_pct": round(shape.top / slide_h, 4),
            "left_pct": round(shape.left / slide_w, 4),
            "width_pct": round(shape.width / slide_w, 4),
            "height_pct": round(shape.height / slide_h, 4),
        }
    except Exception:
        return {}


def _get_theme_el(prs):
    from lxml import etree
    master_part = prs.slide_master.part
    for rel in master_part.rels.values():
        if rel.reltype == _THEME_RT and not rel.is_external:
            return etree.fromstring(rel._target.blob)
    return None


def _resolve_color(el) -> str | None:
    """Extract #RRGGBB from a color container element (srgbClr or sysClr)."""
    if el is None:
        return None
    srgb = el.find("a:srgbClr", _NS)
    if srgb is not None:
        val = srgb.get("val", "")
        return f"#{val.upper()}" if len(val) == 6 else None
    sys_clr = el.find("a:sysClr", _NS)
    if sys_clr is not None:
        val = sys_clr.get("lastClr", "")
        return f"#{val.upper()}" if len(val) == 6 else None
    return None


def _is_dark(hex_color: str) -> bool:
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255 < 0.5


_SCHEME_ACCENT_INDEX = {
    f"accent{i}": i - 1 for i in range(1, 7)  # accent1→0 … accent6→5
}


def _find_brand_accent_index(prs, n_accents: int) -> int:
    """Scan slide XML for the most-referenced scheme accent color.

    Counts both schemeClr (theme references) and srgbClr occurrences that
    match extracted accent_colors, then returns the index of the dominant one.
    This identifies the template's brand color regardless of accent ordering.
    """
    from lxml import etree
    if n_accents == 0:
        return 0

    counts = [0] * n_accents
    _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    for slide in prs.slides:
        try:
            root = slide._element
            for clr_el in root.iter(f"{{{_NS_A}}}schemeClr"):
                val = clr_el.get("val", "")
                if val in _SCHEME_ACCENT_INDEX:
                    idx = _SCHEME_ACCENT_INDEX[val]
                    if idx < n_accents:
                        counts[idx] += 1
        except Exception:
            continue

    return counts.index(max(counts)) if max(counts) > 0 else 0


def _extract_tokens(prs, reusable_slides: dict) -> dict:
    tokens: dict = {}

    theme_el = _get_theme_el(prs)
    if theme_el is not None:
        # Fonts
        major = theme_el.find(".//a:majorFont/a:latin", _NS)
        minor = theme_el.find(".//a:minorFont/a:latin", _NS)
        if major is not None:
            tf = major.get("typeface", "")
            if tf and not tf.startswith("+"):
                tokens["heading_font"] = tf
        if minor is not None:
            tf = minor.get("typeface", "")
            if tf and not tf.startswith("+"):
                tokens["body_font"] = tf

        # Colors
        clr = theme_el.find(".//a:clrScheme", _NS)
        if clr is not None:
            dk1 = _resolve_color(clr.find("a:dk1", _NS))
            lt1 = _resolve_color(clr.find("a:lt1", _NS))
            accents = [_resolve_color(clr.find(f"a:accent{i}", _NS)) for i in range(1, 7)]
            accents = [c for c in accents if c]
            if dk1:
                tokens["dk1_hex"] = dk1
            if lt1:
                tokens["lt1_hex"] = lt1
            if accents:
                tokens["accent_colors"] = accents[:4]
                tokens["brand_accent_index"] = _find_brand_accent_index(prs, len(accents[:4]))
            # bg_is_dark: lt1 is the light background color; if it's actually dark, bg is dark
            if lt1:
                tokens["bg_is_dark"] = _is_dark(lt1)
            elif dk1:
                tokens["bg_is_dark"] = not _is_dark(dk1)

    # Font sizes from slot aggregation (exclude visual_only + decorative_heavy slides)
    sizes: list[int] = []
    max_sz, max_font = -1, None
    font_freq: dict[str, int] = {}
    for sdef in reusable_slides.values():
        if sdef.get("decorative_heavy"):
            continue
        for slot in sdef.get("slots", {}).values():
            if slot.get("visual_only"):
                continue
            sz = slot.get("font_size_pt")
            if sz and 6 <= sz <= 200:
                sizes.append(sz)
                if sz > max_sz:
                    max_sz = sz
                    max_font = slot.get("font_name")
            fn = slot.get("font_name")
            if fn:
                font_freq[fn] = font_freq.get(fn, 0) + 1

    if sizes:
        sizes_s = sorted(sizes)
        n = len(sizes_s)
        tokens["heading_size_pt"] = min(sizes_s[min(int(n * 0.8), n - 1)], 72)
        tokens["body_size_pt"] = max(sizes_s[max(int(n * 0.2), 0)], 8)

    # Fallback font names from slot data if theme didn't provide them
    if "heading_font" not in tokens and max_font:
        tokens["heading_font"] = max_font
    if "body_font" not in tokens and font_freq:
        tokens["body_font"] = max(font_freq, key=font_freq.get)

    return tokens


def _iter_shapes(shapes):
    """Yield all shapes recursively, flattening GROUP containers."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _count_names(shapes) -> dict[str, int]:
    counts: dict[str, int] = {}
    for shape in _iter_shapes(shapes):
        if shape.has_table:
            if any(
                cell.text_frame.text.strip()
                for row in shape.table.rows
                for cell in row.cells
            ):
                counts[shape.name] = counts.get(shape.name, 0) + 1
        elif shape.has_text_frame and shape.text_frame.text.strip():
            counts[shape.name] = counts.get(shape.name, 0) + 1
    return counts


def generate(pptx_path: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    reusable_slides: dict = {}

    n_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        total_count = _count_names(slide.shapes)
        slots: dict = {}
        occurrence: dict[str, int] = {}
        word_counts: list[int] = []  # for decorative_heavy calc (text shapes only)

        for shape in _iter_shapes(slide.shapes):
            sname = shape.name

            # --- TABLE ---
            if shape.has_table:
                rows_data: list[list[str]] = []
                has_text = False
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        ct = cell.text_frame.text.strip()
                        row_data.append(ct)
                        if ct:
                            has_text = True
                    rows_data.append(row_data)
                if not has_text:
                    continue

                n = occurrence.get(sname, 0)
                occurrence[sname] = n + 1
                slot_key = sname if n == 0 else f"{sname}_{n}"

                slot: dict = {
                    "kind": "table",
                    "shape_name": sname,
                    "required": True,
                    "default": rows_data,
                }
                if total_count.get(sname, 0) > 1:
                    slot["nth"] = n
                slot.update(_layout_meta(shape, slide_w, slide_h))
                slots[slot_key] = slot
                continue

            # --- TEXT FRAME ---
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            font_name, font_size = _get_font_info(shape)
            visual_only = _is_visual_only(font_name, text)

            n = occurrence.get(sname, 0)
            occurrence[sname] = n + 1
            slot_key = sname if n == 0 else f"{sname}_{n}"

            paras_with_text = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            kind = "multiline" if len(paras_with_text) > 1 else "text"

            slot = {
                "kind": kind,
                "shape_name": sname,
                "required": True,
                "default": text[:120],
            }
            if visual_only:
                slot["visual_only"] = True
            if total_count.get(sname, 0) > 1:
                slot["nth"] = n
            slot.update(_layout_meta(shape, slide_w, slide_h))
            if font_name:
                slot["font_name"] = font_name
            if font_size:
                slot["font_size_pt"] = font_size
            slots[slot_key] = slot

            word_counts.append((len(text.split()), len(text)))

        if not slots:
            continue

        # Decorative-heavy detection
        # avg_chars < 12 avoids Chinese text false positives (Chinese sentences have no spaces but many chars)
        # "short" shape = few words AND few chars (e.g. "01", "logo")
        n_text = len(word_counts)
        avg_chars = sum(c for _, c in word_counts) / max(n_text, 1)
        short_count = sum(1 for w, c in word_counts if w <= 3 and c <= 8)
        short_ratio = short_count / max(n_text, 1)
        decorative_heavy = (n_text > 10 and avg_chars < 12) or (n_text > 6 and short_ratio > 0.7)

        layout_name = ""
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            pass

        purpose = f"Slide {i + 1}"
        if layout_name:
            purpose += f" — {layout_name}"

        _infer_slot_labels(slots)

        slide_meta: dict = {
            "source_slide_index": i,
            "purpose": purpose,
            "slots": slots,
        }

        intent_tags = _infer_intent_tags(layout_name)
        # Positional fallback: first slide without tags → cover; last → closing
        if not intent_tags:
            if i == 0:
                intent_tags = ["cover"]
            elif i == n_slides - 1:
                intent_tags = ["closing"]
        if intent_tags:
            slide_meta["intent_tags"] = intent_tags

        if decorative_heavy:
            slide_meta["decorative_heavy"] = True
            slide_meta["_meta"] = {
                "text_shape_count": n_text,
                "avg_chars_per_shape": round(avg_chars, 2),
                "short_shape_ratio": round(short_ratio, 2),
                "short_def": "word_count<=3 AND char_count<=8",
            }

        reusable_slides[f"slide_{i}"] = slide_meta

    tokens = _extract_tokens(prs, reusable_slides)

    return {
        "schema_version": "4",
        "template_id": f"auto:{pptx_path.stem}",
        "source_pptx": pptx_path.name,
        "tokens": tokens,
        "reusable_slides": reusable_slides,
        "generated_slides": _GENERATED_SLIDES,
    }


def load_or_generate(pptx_path: Path) -> tuple[dict, Path]:
    """Return (schema, schema_path). Regenerate if schema_version < 4."""
    pptx_path = pptx_path.resolve()
    cache_path = pptx_path.with_suffix(".schema.json")

    if cache_path.exists():
        schema = json.loads(cache_path.read_text(encoding="utf-8"))
        if schema.get("schema_version", "1") >= "4":
            n = len(schema.get("reusable_slides", {}))
            print(f"Schema: loaded cache — {n} slides ({cache_path.name})")
            return schema, cache_path
        print(f"Schema: v{schema.get('schema_version','1')} cache found, regenerating with v4 features…")

    schema = generate(pptx_path)
    cache_path.write_text(json.dumps(schema, ensure_ascii=False, indent=2), encoding="utf-8")
    n = len(schema["reusable_slides"])
    print(f"Schema: auto-generated {n} slides → {cache_path.name}")
    return schema, cache_path
