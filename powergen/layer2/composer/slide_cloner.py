"""
Clone a reusable slide from source_prs to dest_prs, then fill its slots.

No native python-pptx clone API exists — uses lxml deep-copy + OPC part
management directly.
"""
import warnings
from copy import deepcopy

from lxml import etree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml.ns import qn

from .renderers._common import blank_layout
from .slot_resolver import resolve_repeating_field, resolve_slot

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_REMAP_ATTRS = [f"{{{_R_NS}}}embed", f"{{{_R_NS}}}id", f"{{{_R_NS}}}link"]
_SKIP_RELTYPES = {RT.SLIDE_LAYOUT, RT.NOTES_SLIDE, RT.TAGS}


def clone_and_fill(src_prs, dest_prs, slide_def, fill):
    """
    Clone source slide into dest_prs and fill its text/repeating slots.

    Args:
        src_prs:    source Presentation
        dest_prs:   destination Presentation (slide appended in place)
        slide_def:  dict from schema['reusable_slides'][key]
        fill:       {slot_key: value} — str for text/multiline,
                    list-of-dicts for repeating, None for optional_hint

    Returns:
        The newly added pptx.slide.Slide.
    """
    src_slide = src_prs.slides[slide_def["source_slide_index"]]
    new_slide = _clone_slide(src_slide, dest_prs)
    _fill_slots(new_slide, slide_def.get("slots", {}), fill)
    return new_slide


# ---------------------------------------------------------------------------
# Clone
# ---------------------------------------------------------------------------

def _clone_slide(src_slide, dest_prs):
    new_slide = dest_prs.slides.add_slide(blank_layout(dest_prs))

    src_sp_tree = src_slide.shapes._spTree
    new_sp_tree = new_slide.shapes._spTree

    for child in list(new_sp_tree):
        new_sp_tree.remove(child)
    for child in src_sp_tree:
        new_sp_tree.append(deepcopy(child))

    rid_map = {}
    dest_pkg = dest_prs.part.package
    for rid, rel in src_slide.part.rels.items():
        if rel.is_external or rel.reltype in _SKIP_RELTYPES:
            continue
        tgt = rel.target_part
        ext = str(tgt.partname).rsplit(".", 1)[-1]
        new_partname = dest_pkg.next_image_partname(ext)
        new_part = Part(new_partname, tgt.content_type, dest_pkg, tgt.blob)
        new_rid = new_slide.part.relate_to(new_part, rel.reltype)
        rid_map[rid] = new_rid

    if rid_map:
        for el in new_sp_tree.iter():
            for attr in _REMAP_ATTRS:
                val = el.get(attr)
                if val and val in rid_map:
                    el.set(attr, rid_map[val])

    src_cSld = src_slide._element.find(qn("p:cSld"))
    if src_cSld is not None:
        src_bg = src_cSld.find(qn("p:bg"))
        if src_bg is not None:
            new_cSld = new_slide._element.find(qn("p:cSld"))
            existing = new_cSld.find(qn("p:bg"))
            if existing is not None:
                new_cSld.remove(existing)
            new_cSld.insert(0, deepcopy(src_bg))

    return new_slide


# ---------------------------------------------------------------------------
# Slot filling
# ---------------------------------------------------------------------------

def _fill_slots(new_slide, slots_def, fill):
    for slot_key, slot_def in slots_def.items():
        kind = slot_def.get("kind")
        value = fill.get(slot_key) if fill else None
        if value is None:
            value = slot_def.get("default")

        if kind == "repeating":
            instances = (fill.get(slot_key) if fill else None) or []
            _fill_repeating(new_slide, slot_def, instances)
            continue

        if kind == "optional_hint":
            _fill_optional_hint(new_slide, slot_def, value)
            continue

        if value is None:
            continue

        try:
            shape = resolve_slot(new_slide, slot_def)
        except ValueError as exc:
            warnings.warn(f"slot {slot_key!r}: {exc}")
            continue

        if kind == "multiline":
            _set_multiline(shape, str(value))
        elif kind == "image":
            warnings.warn(f"slot {slot_key!r}: image kind deferred in v1, skipping")
        else:
            _set_text(shape, str(value))


def _fill_repeating(new_slide, slot_def, instances):
    max_count = slot_def.get("max_count", len(instances))
    stride_x = float(slot_def.get("stride_x", 0.0))
    stride_y = float(slot_def.get("stride_y", 0.0))
    fields_def = slot_def.get("fields", {})

    for i, instance in enumerate(instances[:max_count]):
        for field_key, field_def in fields_def.items():
            value = instance.get(field_key)
            if value is None:
                value = field_def.get("default")
            if value is None:
                continue
            try:
                shape = resolve_repeating_field(
                    new_slide, field_def, i, stride_x=stride_x, stride_y=stride_y
                )
            except ValueError as exc:
                warnings.warn(f"repeating field {field_key!r}[{i}]: {exc}")
                continue
            if field_def.get("kind") == "multiline":
                _set_multiline(shape, str(value))
            else:
                _set_text(shape, str(value))


def _fill_optional_hint(new_slide, slot_def, value):
    try:
        shape = resolve_slot(new_slide, slot_def)
    except ValueError as exc:
        warnings.warn(f"optional_hint {slot_def.get('shape_name')!r}: {exc}")
        return
    _set_text(shape, str(value) if value is not None else "")


# ---------------------------------------------------------------------------
# Text setters
# ---------------------------------------------------------------------------

def _set_text(shape, text):
    if not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    first_para = paras[0]
    runs = first_para.findall(qn("a:r"))

    for r in runs[1:]:
        first_para.remove(r)
    for p in paras[1:]:
        txBody.remove(p)

    if runs:
        t_el = runs[0].find(qn("a:t"))
        if t_el is None:
            t_el = etree.SubElement(runs[0], qn("a:t"))
        t_el.text = text
    else:
        new_r = etree.SubElement(first_para, qn("a:r"))
        end_rpr = first_para.find(qn("a:endParaRPr"))
        if end_rpr is not None:
            first_para.remove(new_r)
            first_para.insert(list(first_para).index(end_rpr), new_r)
        etree.SubElement(new_r, qn("a:t")).text = text


def _set_multiline(shape, text):
    if not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    first_para = paras[0]

    runs = first_para.findall(qn("a:r"))
    rPr_template = None
    if runs:
        rpr = runs[0].find(qn("a:rPr"))
        if rpr is not None:
            rPr_template = deepcopy(rpr)

    lines = text.split("\n")
    _fill_para(first_para, lines[0], rPr_template)

    for p in paras[1:]:
        txBody.remove(p)
    for line in lines[1:]:
        new_p = deepcopy(first_para)
        txBody.append(new_p)
        _fill_para(new_p, line, rPr_template)


def _fill_para(para_el, line_text, rPr_template):
    for r in para_el.findall(qn("a:r")):
        para_el.remove(r)
    for b in para_el.findall(qn("a:br")):
        para_el.remove(b)

    parts = line_text.split("\v")
    for j, part in enumerate(parts):
        if j > 0:
            br = etree.SubElement(para_el, qn("a:br"))
            if rPr_template is not None:
                br.append(deepcopy(rPr_template))
        new_r = etree.SubElement(para_el, qn("a:r"))
        if rPr_template is not None:
            new_r.append(deepcopy(rPr_template))
        etree.SubElement(new_r, qn("a:t")).text = part

    end_rpr = para_el.find(qn("a:endParaRPr"))
    if end_rpr is not None:
        para_el.remove(end_rpr)
        para_el.append(end_rpr)
