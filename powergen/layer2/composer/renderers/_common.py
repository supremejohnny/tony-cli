"""Shared utilities for all content-type renderers."""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def rgb(hex_str):
    """Parse #RRGGBB → RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def blank_layout(prs):
    """Return the Blank slide layout, or the layout with fewest placeholders."""
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout
    return min(prs.slide_layouts, key=lambda l: len(list(l.placeholders)))


def slide_dims(prs):
    """Return (width_in, height_in) as floats."""
    from pptx.util import Emu
    return Emu(prs.slide_width).inches, Emu(prs.slide_height).inches


def add_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    """Add a solid-filled rectangle (no border by default)."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    border_color = fill_hex if line_hex is None else line_hex
    shape.line.color.rgb = rgb(border_color)
    return shape


def add_text(
    slide, text, left, top, width, height,
    font_size=12, bold=False, color_hex=None,
    align=PP_ALIGN.LEFT, word_wrap=True, italic=False,
):
    """Add a textbox and return the shape."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color_hex:
        run.font.color.rgb = rgb(color_hex)
    return box


def add_text_multiline(
    slide, lines, left, top, width, height,
    font_size=12, bold=False, color_hex=None,
    align=PP_ALIGN.LEFT, word_wrap=True,
):
    """Add a textbox with multiple paragraphs."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = word_wrap
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color_hex:
            run.font.color.rgb = rgb(color_hex)
    return box


def primary(tokens):
    return tokens.get("primary_color", "#1F2937")


def accent(tokens, index=0):
    accents = tokens.get("accent_colors", ["#7578EC", "#F7B802", "#F18703", "#F35B06"])
    return accents[index % len(accents)] if accents else "#7578EC"


def title_size(tokens):
    r = tokens.get("title_font", {}).get("size_pt_range", [28, 40])
    return r[1]  # use max


def body_size(tokens):
    r = tokens.get("body_font", {}).get("size_pt_range", [12, 16])
    return r[1]
