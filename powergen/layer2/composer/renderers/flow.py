"""
Flow renderer — horizontal A → B → C arrow chain.

Fill format:
    {
        "title": "Slide title",
        "items": ["Step A", "Step B", "Step C"]
    }

Max 5 items. Items are placed in rounded rectangles connected by arrows.
"""
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ._common import (
    brand_accent, add_rect, add_text, blank_layout,
    body_size, primary, rgb, slide_dims, title_size,
)

MAX_ITEMS = 5
BAR_W = 0.28
MARGIN_L = BAR_W + 0.35
MARGIN_T = 0.30
MARGIN_R = 0.35
ARROW_W = 0.35
BOX_H = 1.4
ARROW_COLOR = "#9CA3AF"


def render(prs, fill, tokens):
    """Add one flow slide to prs."""
    slide_title = fill.get("title", "")
    items = fill.get("items", [])[:MAX_ITEMS]

    layout = blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    _render_flow_slide(slide, prs, slide_title, items, tokens)


def _render_flow_slide(slide, prs, slide_title, items, tokens):
    w, h = slide_dims(prs)
    content_w = w - MARGIN_L - MARGIN_R
    pri = primary(tokens)
    t_size = min(title_size(tokens), 36)
    b_size = body_size(tokens)
    n = len(items)

    # Left accent bar
    add_rect(slide, 0, 0, BAR_W, h, pri)

    # Slide title
    add_text(
        slide, slide_title,
        left=MARGIN_L, top=MARGIN_T, width=content_w, height=0.8,
        font_size=t_size, bold=True, color_hex=pri,
    )

    if not items:
        return

    total_arrow = (n - 1) * ARROW_W
    box_w = (content_w - total_arrow) / n
    box_top = (h - BOX_H) / 2 + 0.2  # vertically centered, slight offset for title

    for i, text in enumerate(items):
        box_left = MARGIN_L + i * (box_w + ARROW_W)
        acc = brand_accent(tokens, i)

        # Box background
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(box_left), Inches(box_top), Inches(box_w), Inches(BOX_H),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(acc)
        box.line.color.rgb = rgb(acc)

        # Number badge
        badge_size = 0.3
        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(box_left + box_w / 2 - badge_size / 2),
            Inches(box_top + 0.1),
            Inches(badge_size), Inches(badge_size),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb("#FFFFFF")
        badge.line.color.rgb = rgb(acc)

        num_tf = badge.text_frame
        num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_run = num_tf.paragraphs[0].add_run()
        num_run.text = str(i + 1)
        num_run.font.size = Pt(9)
        num_run.font.bold = True
        num_run.font.color.rgb = rgb(acc)

        # Item text
        add_text(
            slide, text,
            left=box_left + 0.1, top=box_top + 0.5,
            width=box_w - 0.2, height=BOX_H - 0.6,
            font_size=b_size, bold=False, color_hex="#FFFFFF",
            align=PP_ALIGN.CENTER, word_wrap=True,
        )

        # Arrow (except after last item)
        if i < n - 1:
            arrow_left = box_left + box_w
            arrow_cx = arrow_left + ARROW_W / 2
            arrow_cy = box_top + BOX_H / 2
            add_text(
                slide, "→",
                left=arrow_left, top=arrow_cy - 0.2,
                width=ARROW_W, height=0.4,
                font_size=18, color_hex=ARROW_COLOR,
                align=PP_ALIGN.CENTER, word_wrap=False,
            )
