"""
Validate a template.schema.json against its source .pptx.

Usage:  python validate.py <path/to/template.schema.json>

Checks performed:
  1. JSON parses, schema_version is "1"
  2. source_pptx file exists
  3. Every reusable_slides[*].source_slide_index is in range
  4. Every slot locator (shape_name [+ nth] [+ near]) resolves to exactly one shape
  5. Every accent_color and primary_color is valid hex #RRGGBB
  6. Every generated_slides key is in the registered content_type set
  7. compose_hints.section_intro_pairs reference existing reusable_slides keys

Strips _-prefixed comment fields before validation.
"""
import json
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

REGISTERED_CONTENT_TYPES = {"bullet", "numbered", "card", "flow", "text_block"}
HEX_RE = re.compile(r"^#[0-9A-Fa-f]{6}$")


def strip_comments(obj):
    if isinstance(obj, dict):
        return {k: strip_comments(v) for k, v in obj.items() if not k.startswith("_")}
    if isinstance(obj, list):
        return [strip_comments(v) for v in obj]
    return obj


def inches(v):
    return Emu(v).inches if v is not None else None


def resolve_locator(slide, shape_name, nth=None, near=None):
    """Return list of shapes matching the locator. Should be exactly 1."""
    candidates = [s for s in slide.shapes if s.name == shape_name]
    if not candidates:
        return []
    if nth is not None:
        return [candidates[nth]] if 0 <= nth < len(candidates) else []
    if near is not None:
        target_t, target_l = near.get("top"), near.get("left")
        def dist(s):
            return abs(inches(s.top) - target_t) + abs(inches(s.left) - target_l)
        candidates.sort(key=dist)
        return [candidates[0]]
    return candidates  # caller checks len == 1


def validate(schema_path):
    errors = []
    schema_path = Path(schema_path)
    raw = json.loads(schema_path.read_text())
    schema = strip_comments(raw)

    # 1. version
    if schema.get("schema_version") != "1":
        errors.append(f"schema_version must be '1', got {schema.get('schema_version')!r}")

    # 2. source pptx exists
    src = schema_path.parent / schema["source_pptx"]
    if not src.exists():
        errors.append(f"source_pptx not found: {src}")
        return errors  # can't continue without the pptx

    pres = Presentation(src)
    n_slides = len(pres.slides)

    # 5. tokens hex
    tokens = schema.get("tokens", {})
    primary = tokens.get("primary_color")
    if primary and not HEX_RE.match(primary):
        errors.append(f"tokens.primary_color invalid hex: {primary!r}")
    for c in tokens.get("accent_colors", []):
        if not HEX_RE.match(c):
            errors.append(f"tokens.accent_colors contains invalid hex: {c!r}")

    # 3 & 4. reusable slides + slot locators
    reusable = schema.get("reusable_slides", {})
    for skey, sdef in reusable.items():
        idx = sdef["source_slide_index"]
        if not (0 <= idx < n_slides):
            errors.append(f"reusable_slides[{skey}].source_slide_index={idx} out of range [0,{n_slides})")
            continue
        slide = pres.slides[idx]
        for slot_key, slot in sdef.get("slots", {}).items():
            if slot["kind"] == "repeating":
                anchor = slot.get("anchor_shape")
                if anchor and not any(s.name == anchor for s in slide.shapes):
                    errors.append(f"[{skey}.{slot_key}] anchor_shape={anchor!r} not found on slide {idx}")
                for fkey, fdef in slot.get("fields", {}).items():
                    matches = [s for s in slide.shapes if s.name == fdef["shape_name"]]
                    if not matches:
                        errors.append(f"[{skey}.{slot_key}.{fkey}] shape_name={fdef['shape_name']!r} not found on slide {idx}")
            else:
                matches = resolve_locator(slide, slot["shape_name"], slot.get("nth"), slot.get("near"))
                if len(matches) == 0:
                    errors.append(f"[{skey}.{slot_key}] locator did not resolve on slide {idx}: name={slot['shape_name']!r} nth={slot.get('nth')}")
                elif len(matches) > 1 and slot.get("nth") is None and slot.get("near") is None:
                    errors.append(f"[{skey}.{slot_key}] ambiguous: {len(matches)} shapes named {slot['shape_name']!r} on slide {idx}, add nth or near")

    # 6. content_type registration
    for ct in schema.get("generated_slides", {}):
        if ct not in REGISTERED_CONTENT_TYPES:
            errors.append(f"generated_slides[{ct!r}] not in registered set {REGISTERED_CONTENT_TYPES}")

    # 7. compose_hints references
    for pair in schema.get("compose_hints", {}).get("section_intro_pairs", []):
        for role in ("intro", "divider"):
            ref = pair.get(role)
            if ref and ref not in reusable:
                errors.append(f"compose_hints.section_intro_pairs references unknown reusable {ref!r}")

    return errors


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python validate.py <path/to/template.schema.json>")
        sys.exit(1)
    errs = validate(sys.argv[1])
    if errs:
        print(f"FAIL: {len(errs)} error(s)")
        for e in errs:
            print(f"  - {e}")
        sys.exit(1)
    print("OK")
