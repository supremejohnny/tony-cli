"""
Smoke test for slide_cloner: clone all reusable slides and fill with sample data.

Usage:
    python -X utf8 -m powergen.layer2.scripts.test_cloner

Output: test/layer2_cloner_test.pptx
Verify: python -X utf8 -m powergen.layer2.scripts.inspect_pptx test/layer2_cloner_test.pptx
"""
from pathlib import Path

from pptx import Presentation

from powergen.layer2.composer.schema_loader import load
from powergen.layer2.composer.slide_cloner import clone_and_fill

SCHEMA_PATH = Path("powergen/layer2/schemas/test_template.schema.json")
OUT_PATH = Path("test/layer2_cloner_test.pptx")

FILLS = {
    "cover": {
        "university": "University of Toronto",
        "student_info": "路觅学生：张伟",
        "advisor_info": "路觅导师：李明",
        "semester_label": "Fall 2026",
    },
    "mentor_intro": {
        "mentor_name": "Dr. Smith",
        "mentor_credentials": "MIT\nPhD Computer Science",
        "mentor_bio": "Ten years of teaching experience at MIT and Stanford.",
        "mentor_achievements": "500+ students mentored\v95% grad school admission rate",
    },
    "section_intro": {
        "section_title": "课程规划",
        "section_subtitle": "Course Planning Overview",
    },
    "numbered_card_grid_v": {
        "section_title_zh": "选课红线",
        "section_title_en": "Course Selection Red Lines",
        "cards": [
            {
                "number": "01",
                "title": "最低学分要求",
                "body": "每学期至少修读12学分课程。",
                "footer": "低于此要求将影响F-1签证身份。",
            },
            {
                "number": "02",
                "title": "最高学分限制",
                "body": "每学期最多修读20学分课程。",
                "footer": "超出需向学院申请豁免。",
            },
        ],
    },
    "numbered_card_grid_h": {
        "section_title_zh": "后果分析",
        "section_title_en": "Consequence Analysis",
        "cards": [
            {"number": "01", "category": "经济", "body": "重修课程需缴纳全额学费。"},
            {"number": "02", "category": "时间", "body": "延误毕业进程至少一学期。"},
            {"number": "03", "category": "签证", "body": "影响F-1身份合规性。"},
            {"number": "04", "category": "专业", "body": "阻碍进入目标专业通道。"},
        ],
    },
    "section_divider": {
        "number": "02.",
        "title_zh": "课程规划",
        "subtitle_en": "Course Planning Strategy",
    },
    "closing_timeline": {
        "title": "时间线",
        "timeline_notes": None,
    },
}


def main():
    schema = load(SCHEMA_PATH)
    src_prs = Presentation("test/test.pptx")

    dest_prs = Presentation()
    dest_prs.slide_width = src_prs.slide_width
    dest_prs.slide_height = src_prs.slide_height

    for key, slide_def in schema["reusable_slides"].items():
        fill = FILLS.get(key, {})
        slide = clone_and_fill(src_prs, dest_prs, slide_def, fill)
        print(f"cloned: {key} — {len(list(slide.shapes))} shapes")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    dest_prs.save(str(OUT_PATH))
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
