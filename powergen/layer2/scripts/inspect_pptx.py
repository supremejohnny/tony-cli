"""
Dump shape inventory from a .pptx file.

Usage:
    python -m powergen.layer2.scripts.inspect_pptx <path/to/file.pptx>
    python -m powergen.layer2.scripts.inspect_pptx <file.pptx> --slide 3

Output: markdown table per slide — use as Step 1 of the schema authoring procedure in SKILL.md.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def _in(emu):
    return f"{Emu(emu).inches:.2f}" if emu is not None else "-"


def _text_preview(shape, limit=60):
    try:
        if shape.has_text_frame:
            text = shape.text_frame.text.replace("\n", "\\n").replace("\r", "")
            return (text[:limit] + "…") if len(text) > limit else text
    except Exception:
        pass
    return ""


def _shape_rows(shapes, depth=0):
    prefix = "  " * depth
    rows = []
    for i, shape in enumerate(shapes):
        rows.append(
            f"| {prefix}{i} | {shape.name} | {shape.shape_type} "
            f"| {_in(shape.top)} | {_in(shape.left)} "
            f"| {_in(shape.width)}×{_in(shape.height)} "
            f"| {_text_preview(shape)} |"
        )
        if hasattr(shape, "shapes"):  # group
            rows.extend(_shape_rows(shape.shapes, depth + 1))
    return rows


def inspect(pptx_path, slide_filter=None):
    prs = Presentation(str(pptx_path))
    lines = [
        f"# Shape inventory: {Path(pptx_path).name}",
        f"Slides: {len(prs.slides)}",
    ]
    for i, slide in enumerate(prs.slides):
        if slide_filter is not None and i != slide_filter:
            continue
        layout = slide.slide_layout.name if slide.slide_layout else "(none)"
        lines.append(f"\n## Slide {i} — layout: {layout!r} — {len(slide.shapes)} shape(s)")
        lines.append("| # | name | type | top_in | left_in | size_in | text_preview |")
        lines.append("|---|------|------|--------|---------|---------|-------------|")
        lines.extend(_shape_rows(slide.shapes))
    return "\n".join(lines)


if __name__ == "__main__":
    args = sys.argv[1:]
    if not args or args[0] in ("-h", "--help"):
        print(__doc__)
        sys.exit(0)

    pptx_path = args[0]
    slide_filter = None
    if "--slide" in args:
        idx = args.index("--slide")
        slide_filter = int(args[idx + 1])

    if not Path(pptx_path).exists():
        print(f"Error: file not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    print(inspect(pptx_path, slide_filter))
